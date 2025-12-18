import streamlit as st
import pandas as pd
import sqlite3
from datetime import datetime, timedelta
import plotly.express as px
import plotly.graph_objects as go
from plotly.subplots import make_subplots
import json
from io import BytesIO
import tempfile
import os

# Page configuration
st.set_page_config(
    page_title="VCU Simulation Center Work Tracker",
    page_icon="🏥",
    layout="wide",
    initial_sidebar_state="expanded"
)

# Room numbers constant
ROOM_NUMBERS = ["9-203", "9-205", "9-208", "9-209", "9-210", "9-211", "9-215", "9-217"]

def get_active_rooms():
    """Get list of active room numbers"""
    rooms_df = get_room_numbers(active_only=True)
    if not rooms_df.empty:
        return sorted(rooms_df['room_number'].tolist())
    return ROOM_NUMBERS  # Fallback to constant if DB empty

# Database path
DB_PATH = 'work_tracker.db'

# Database setup
def init_db():
    conn = sqlite3.connect(DB_PATH)
    c = conn.cursor()
    
    # Main activities table
    c.execute('''
        CREATE TABLE IF NOT EXISTS activities (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            date DATE NOT NULL,
            activity_type TEXT NOT NULL,
            hours REAL NOT NULL,
            students_trained INTEGER DEFAULT 0,
            personnel TEXT,
            equipment TEXT,
            course TEXT,
            room_number TEXT,
            time_start TEXT,
            time_end TEXT,
            turn_in INTEGER DEFAULT 0,
            received INTEGER DEFAULT 0,
            notes TEXT,
            created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
        )
    ''')

    # Migration: Add room columns if they don't exist
    c.execute("PRAGMA table_info(activities)")
    columns = [col[1] for col in c.fetchall()]
    
    if 'room_number' not in columns:
        c.execute("ALTER TABLE activities ADD COLUMN room_number TEXT")
    if 'time_start' not in columns:
        c.execute("ALTER TABLE activities ADD COLUMN time_start TEXT")
    if 'time_end' not in columns:
        c.execute("ALTER TABLE activities ADD COLUMN time_end TEXT")
    
    # Personnel table
    c.execute('''
        CREATE TABLE IF NOT EXISTS personnel (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            name TEXT NOT NULL UNIQUE,
            role TEXT,
            active BOOLEAN DEFAULT 1,
            created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
        )
    ''')
    
    # Equipment table
    c.execute('''
        CREATE TABLE IF NOT EXISTS equipment (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            name TEXT NOT NULL UNIQUE,
            status TEXT DEFAULT 'Operational',
            last_maintenance DATE,
            notes TEXT,
            active BOOLEAN DEFAULT 1,
            created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
        )
    ''')
    
    # Courses table (simplified)
    c.execute('''
        CREATE TABLE IF NOT EXISTS courses (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            name TEXT NOT NULL UNIQUE,
            description TEXT,
            active BOOLEAN DEFAULT 1,
            created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
        )
    ''')    
    # Activity types table
    c.execute('''
        CREATE TABLE IF NOT EXISTS activity_types (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            name TEXT NOT NULL UNIQUE,
            description TEXT,
            active BOOLEAN DEFAULT 1,
            created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
        )
    ''')

    # Migration: Add room columns if they don't exist
    c.execute("PRAGMA table_info(activities)")
    columns = [col[1] for col in c.fetchall()]
    
    if 'room_number' not in columns:
        c.execute("ALTER TABLE activities ADD COLUMN room_number TEXT")
    if 'time_start' not in columns:
        c.execute("ALTER TABLE activities ADD COLUMN time_start TEXT")
    if 'time_end' not in columns:
        c.execute("ALTER TABLE activities ADD COLUMN time_end TEXT")
    
    # Incidents table
    c.execute('''
        CREATE TABLE IF NOT EXISTS incidents (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            date DATE NOT NULL,
            incident_type TEXT NOT NULL,
            equipment TEXT,
            severity TEXT,
            description TEXT,
            resolution TEXT,
            resolved BOOLEAN DEFAULT 0,
            created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
        )
    ''')

    # Migration: Add room columns if they don't exist
    c.execute("PRAGMA table_info(activities)")
    columns = [col[1] for col in c.fetchall()]
    
    if 'room_number' not in columns:
        c.execute("ALTER TABLE activities ADD COLUMN room_number TEXT")
    if 'time_start' not in columns:
        c.execute("ALTER TABLE activities ADD COLUMN time_start TEXT")
    if 'time_end' not in columns:
        c.execute("ALTER TABLE activities ADD COLUMN time_end TEXT")
    
    # Goals table
    c.execute('''
        CREATE TABLE IF NOT EXISTS goals (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            goal_type TEXT NOT NULL,
            target_value REAL NOT NULL,
            current_value REAL DEFAULT 0,
            period TEXT NOT NULL,
            status TEXT DEFAULT 'In Progress',
            created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
        )
    ''')
    
    # Cancellations table (NEW!)
    c.execute('''
        CREATE TABLE IF NOT EXISTS cancellations (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            activity_id INTEGER,
            date DATE NOT NULL,
            course TEXT,
            scheduled_time TEXT,
            scheduled_duration REAL,
            reason TEXT NOT NULL,
            notes TEXT,
            impacted_students INTEGER DEFAULT 0,
            rescheduled BOOLEAN DEFAULT 0,
            reschedule_date DATE,
            tech_time_spent REAL DEFAULT 0,
            activity_type TEXT,
            personnel TEXT,
            equipment TEXT,
            room_number TEXT,
            created_by TEXT,
            created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
        )
    ''')
    
    # Migration: Add new columns to cancellations table if they don't exist
    c.execute("PRAGMA table_info(cancellations)")
    cancel_columns = [col[1] for col in c.fetchall()]
    
    if 'activity_id' not in cancel_columns:
        c.execute("ALTER TABLE cancellations ADD COLUMN activity_id INTEGER")
    if 'tech_time_spent' not in cancel_columns:
        c.execute("ALTER TABLE cancellations ADD COLUMN tech_time_spent REAL DEFAULT 0")
    if 'activity_type' not in cancel_columns:
        c.execute("ALTER TABLE cancellations ADD COLUMN activity_type TEXT")
    if 'personnel' not in cancel_columns:
        c.execute("ALTER TABLE cancellations ADD COLUMN personnel TEXT")
    if 'equipment' not in cancel_columns:
        c.execute("ALTER TABLE cancellations ADD COLUMN equipment TEXT")
    if 'room_number' not in cancel_columns:
        c.execute("ALTER TABLE cancellations ADD COLUMN room_number TEXT")
    
    # Time Off tracker table (NEW!)
    c.execute('''
        CREATE TABLE IF NOT EXISTS time_off (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            personnel TEXT NOT NULL,
            start_date DATE NOT NULL,
            end_date DATE NOT NULL,
            time_off_type TEXT NOT NULL,
            hours REAL NOT NULL,
            status TEXT DEFAULT 'Approved',
            notes TEXT,
            created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
        )
    ''')
    
    # Leave Types table (NEW!)
    c.execute('''
        CREATE TABLE IF NOT EXISTS leave_types (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            leave_type_name TEXT NOT NULL UNIQUE,
            default_annual_hours REAL DEFAULT 0,
            active BOOLEAN DEFAULT 1,
            created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
        )
    ''')
    
    # Initialize default leave types if table is empty
    c.execute("SELECT COUNT(*) FROM leave_types")
    if c.fetchone()[0] == 0:
        default_types = [
            ('Annual Leave', 160),
            ('Community Service Leave', 40),
            ('Enhanced Community Service Leave', 80),
            ('Military Leave', 0),
            ('Personal Leave', 24),
            ('Sick Leave', 96),
            ('University Leave', 120),
        ]
        c.executemany("INSERT INTO leave_types (leave_type_name, default_annual_hours) VALUES (?, ?)", default_types)
    
    # Leave Accruals table (NEW!)
    c.execute('''
        CREATE TABLE IF NOT EXISTS leave_accruals (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            personnel TEXT NOT NULL,
            leave_type TEXT NOT NULL,
            hours_available REAL DEFAULT 0,
            hours_used REAL DEFAULT 0,
            last_updated TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
            UNIQUE(personnel, leave_type)
        )
    ''')
    
    # Room Numbers table (NEW!)
    c.execute('''
        CREATE TABLE IF NOT EXISTS room_numbers (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            room_number TEXT NOT NULL UNIQUE,
            active BOOLEAN DEFAULT 1,
            notes TEXT,
            created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
        )
    ''')
    
    # Initialize default rooms if table is empty
    c.execute("SELECT COUNT(*) FROM room_numbers")
    if c.fetchone()[0] == 0:
        default_rooms = [(room,) for room in ROOM_NUMBERS]
        c.executemany("INSERT INTO room_numbers (room_number) VALUES (?)", default_rooms)
    
    conn.commit()
    conn.close()

#############################################
# CANCELLATION TRACKING FUNCTIONS (NEW!)
#############################################

def add_cancellation(date, course, scheduled_time, scheduled_duration, reason, notes, impacted_students, rescheduled, reschedule_date, created_by, activity_id=None, tech_time_spent=0, activity_type=None, personnel=None, equipment=None, room_number=None):
    """Add a course cancellation record"""
    conn = sqlite3.connect(DB_PATH)
    c = conn.cursor()
    c.execute('''
        INSERT INTO cancellations (activity_id, date, course, scheduled_time, scheduled_duration, reason, notes, 
                                   impacted_students, rescheduled, reschedule_date, tech_time_spent,
                                   activity_type, personnel, equipment, room_number, created_by)
        VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?)
    ''', (activity_id, date, course, scheduled_time, scheduled_duration, reason, notes, 
          impacted_students, rescheduled, reschedule_date, tech_time_spent,
          activity_type, personnel, equipment, room_number, created_by))
    conn.commit()
    conn.close()

def cancel_existing_activity(activity_id, reason, notes, tech_time_spent, rescheduled, reschedule_date, created_by):
    """Cancel an existing activity and create cancellation record"""
    conn = sqlite3.connect(DB_PATH)
    c = conn.cursor()
    
    # Get activity details
    c.execute("SELECT * FROM activities WHERE id = ?", (activity_id,))
    activity = c.fetchone()
    
    if activity:
        # Column indices: 0=id, 1=date, 2=activity_type, 3=hours, 4=students_trained, 5=personnel, 6=equipment, 7=course, 8=room_number, 9=time_start, 10=time_end, 11=turn_in, 12=received, 13=notes
        
        # Create cancellation record
        c.execute('''
            INSERT INTO cancellations (activity_id, date, course, scheduled_time, scheduled_duration, 
                                       reason, notes, impacted_students, rescheduled, reschedule_date, 
                                       tech_time_spent, activity_type, personnel, equipment, room_number, created_by)
            VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?)
        ''', (activity_id, activity[1], activity[7], activity[9], activity[3], notes, activity[4], 
              rescheduled, reschedule_date, tech_time_spent, activity[2], activity[5], activity[6], activity[8], created_by))
        
        # Mark activity as cancelled in notes
        old_notes = activity[13] if activity[13] else ""
        new_notes = f"[CANCELLED: {reason}] {old_notes}"
        c.execute("UPDATE activities SET notes = ? WHERE id = ?", (new_notes, activity_id))
        
        conn.commit()
        conn.close()
        return True
    
    conn.close()
    return False

def get_active_activities_for_cancellation(start_date, end_date):
    """Get activities that can be cancelled (not already cancelled)"""
    conn = sqlite3.connect(DB_PATH)
    query = """
        SELECT * FROM activities
        WHERE date BETWEEN ? AND ?
        AND (notes IS NULL OR notes NOT LIKE '%[CANCELLED:%')
        ORDER BY date DESC, time_start DESC
    """
    df = pd.read_sql_query(query, conn, params=(start_date, end_date))
    conn.close()
    return df

def get_cancellations(start_date=None, end_date=None):
    """Get cancellation records, optionally filtered by date range"""
    conn = sqlite3.connect(DB_PATH)
    if start_date and end_date:
        query = "SELECT * FROM cancellations WHERE date BETWEEN ? AND ? ORDER BY date DESC"
        df = pd.read_sql_query(query, conn, params=(start_date, end_date))
    else:
        df = pd.read_sql_query("SELECT * FROM cancellations ORDER BY date DESC", conn)
    conn.close()
    return df

def update_cancellation(cancellation_id, **kwargs):
    """Update a cancellation record"""
    conn = sqlite3.connect(DB_PATH)
    c = conn.cursor()
    
    set_clause = ", ".join([f"{k} = ?" for k in kwargs.keys()])
    values = list(kwargs.values()) + [cancellation_id]
    
    c.execute(f"UPDATE cancellations SET {set_clause} WHERE id = ?", values)
    conn.commit()
    conn.close()

def delete_cancellation(cancellation_id):
    """Delete a cancellation record"""
    conn = sqlite3.connect(DB_PATH)
    c = conn.cursor()
    c.execute("DELETE FROM cancellations WHERE id = ?", (cancellation_id,))
    conn.commit()
    conn.close()

#############################################
# TIME OFF TRACKING FUNCTIONS (NEW!)
#############################################

def add_time_off(personnel, start_date, end_date, time_off_type, hours, status, notes):
    """Add a time off record and auto-deduct from leave accruals"""
    # First, check if there are sufficient leave hours and deduct
    success, remaining = deduct_leave_hours(personnel, time_off_type, hours)
    
    # Add time off record regardless (for audit trail)
    conn = sqlite3.connect(DB_PATH)
    c = conn.cursor()
    c.execute('''
        INSERT INTO time_off (personnel, start_date, end_date, time_off_type, hours, status, notes)
        VALUES (?, ?, ?, ?, ?, ?, ?)
    ''', (personnel, start_date, end_date, time_off_type, hours, status, notes))
    conn.commit()
    conn.close()
    
    return success, remaining

def get_time_off(start_date=None, end_date=None, personnel=None):
    """Get time off records, optionally filtered"""
    conn = sqlite3.connect(DB_PATH)
    
    query = "SELECT * FROM time_off WHERE 1=1"
    params = []
    
    if start_date and end_date:
        query += " AND start_date <= ? AND end_date >= ?"
        params.extend([end_date, start_date])
    
    if personnel:
        query += " AND personnel = ?"
        params.append(personnel)
    
    query += " ORDER BY start_date DESC"
    
    if params:
        df = pd.read_sql_query(query, conn, params=params)
    else:
        df = pd.read_sql_query(query, conn)
    
    conn.close()
    return df

def update_time_off(time_off_id, **kwargs):
    """Update a time off record and adjust leave balances if leave type changes"""
    conn = sqlite3.connect(DB_PATH)
    c = conn.cursor()
    
    # Get original record to check if leave type is changing
    c.execute("SELECT personnel, time_off_type, hours FROM time_off WHERE id = ?", (time_off_id,))
    original = c.fetchone()
    
    if original and 'time_off_type' in kwargs:
        personnel, old_leave_type, hours = original
        new_leave_type = kwargs['time_off_type']
        
        # If leave type is changing, adjust balances
        if old_leave_type != new_leave_type:
            # Add hours back to old leave type (undo the deduction)
            c.execute("""
                UPDATE leave_accruals 
                SET hours_available = hours_available + ?,
                    hours_used = hours_used - ?,
                    last_updated = CURRENT_TIMESTAMP
                WHERE personnel = ? AND leave_type = ?
            """, (hours, hours, personnel, old_leave_type))
            
            # Deduct hours from new leave type
            c.execute("""
                UPDATE leave_accruals 
                SET hours_available = hours_available - ?,
                    hours_used = hours_used + ?,
                    last_updated = CURRENT_TIMESTAMP
                WHERE personnel = ? AND leave_type = ?
            """, (hours, hours, personnel, new_leave_type))
    
    # Update the time_off record
    set_clause = ", ".join([f"{k} = ?" for k in kwargs.keys()])
    values = list(kwargs.values()) + [time_off_id]
    
    c.execute(f"UPDATE time_off SET {set_clause} WHERE id = ?", values)
    conn.commit()
    conn.close()

def delete_time_off(time_off_id):
    """Delete a time off record and restore leave balance"""
    conn = sqlite3.connect(DB_PATH)
    c = conn.cursor()
    
    # Get the record details before deleting
    c.execute("SELECT personnel, time_off_type, hours FROM time_off WHERE id = ?", (time_off_id,))
    record = c.fetchone()
    
    if record:
        personnel, leave_type, hours = record
        
        # Restore the hours to the leave balance
        c.execute("""
            UPDATE leave_accruals 
            SET hours_available = hours_available + ?,
                hours_used = hours_used - ?,
                last_updated = CURRENT_TIMESTAMP
            WHERE personnel = ? AND leave_type = ?
        """, (hours, hours, personnel, leave_type))
    
    # Delete the record
    c.execute("DELETE FROM time_off WHERE id = ?", (time_off_id,))
    conn.commit()
    conn.close()

def get_time_off_summary(year=None):
    """Get summary of time off by personnel"""
    conn = sqlite3.connect(DB_PATH)
    
    if year:
        query = """
            SELECT personnel, 
                   time_off_type,
                   SUM(hours) as total_hours,
                   COUNT(*) as occurrences
            FROM time_off
            WHERE strftime('%Y', start_date) = ?
            GROUP BY personnel, time_off_type
            ORDER BY personnel, time_off_type
        """
        df = pd.read_sql_query(query, conn, params=(str(year),))
    else:
        query = """
            SELECT personnel, 
                   time_off_type,
                   SUM(hours) as total_hours,
                   COUNT(*) as occurrences
            FROM time_off
            GROUP BY personnel, time_off_type
            ORDER BY personnel, time_off_type
        """
        df = pd.read_sql_query(query, conn)
    
    conn.close()
    return df

#############################################
# LEAVE ACCRUAL MANAGEMENT FUNCTIONS (NEW!)
#############################################

def get_leave_types(active_only=True):
    """Get all leave types"""
    conn = sqlite3.connect(DB_PATH)
    if active_only:
        df = pd.read_sql_query("SELECT * FROM leave_types WHERE active = 1 ORDER BY leave_type_name", conn)
    else:
        df = pd.read_sql_query("SELECT * FROM leave_types ORDER BY leave_type_name", conn)
    conn.close()
    return df

def add_leave_type(leave_type_name, default_annual_hours):
    """Add a new leave type"""
    conn = sqlite3.connect(DB_PATH)
    c = conn.cursor()
    try:
        c.execute("INSERT INTO leave_types (leave_type_name, default_annual_hours) VALUES (?, ?)", 
                  (leave_type_name, default_annual_hours))
        conn.commit()
        conn.close()
        return True
    except sqlite3.IntegrityError:
        conn.close()
        return False

def update_leave_type(leave_type_id, **kwargs):
    """Update a leave type"""
    conn = sqlite3.connect(DB_PATH)
    c = conn.cursor()
    set_clause = ", ".join([f"{k} = ?" for k in kwargs.keys()])
    values = list(kwargs.values()) + [leave_type_id]
    c.execute(f"UPDATE leave_types SET {set_clause} WHERE id = ?", values)
    conn.commit()
    conn.close()

def delete_leave_type(leave_type_id):
    """Deactivate a leave type (don't delete - preserve history)"""
    conn = sqlite3.connect(DB_PATH)
    c = conn.cursor()
    c.execute("UPDATE leave_types SET active = 0 WHERE id = ?", (leave_type_id,))
    conn.commit()
    conn.close()

def get_leave_accruals(personnel=None):
    """Get leave accruals for a person or all personnel"""
    conn = sqlite3.connect(DB_PATH)
    if personnel:
        df = pd.read_sql_query(
            "SELECT * FROM leave_accruals WHERE personnel = ? ORDER BY leave_type", 
            conn, params=(personnel,))
    else:
        df = pd.read_sql_query("SELECT * FROM leave_accruals ORDER BY personnel, leave_type", conn)
    conn.close()
    return df

def initialize_leave_accruals(personnel):
    """Initialize leave accruals for a new person - starts at 0 hours for full control"""
    conn = sqlite3.connect(DB_PATH)
    c = conn.cursor()
    
    # Get active leave types
    c.execute("SELECT leave_type_name FROM leave_types WHERE active = 1")
    leave_types = c.fetchall()
    
    # Insert accruals for each leave type starting at 0 hours
    for (leave_type,) in leave_types:
        c.execute('''
            INSERT OR IGNORE INTO leave_accruals (personnel, leave_type, hours_available, hours_used)
            VALUES (?, ?, 0, 0)
        ''', (personnel, leave_type))
    
    conn.commit()
    conn.close()

def add_accrual_hours(personnel, leave_type, hours, reason=None):
    """Add accrual hours to a person's leave balance"""
    conn = sqlite3.connect(DB_PATH)
    c = conn.cursor()
    
    # Check if accrual exists
    c.execute("SELECT hours_available FROM leave_accruals WHERE personnel = ? AND leave_type = ?", 
              (personnel, leave_type))
    result = c.fetchone()
    
    if result:
        new_balance = result[0] + hours
        c.execute("""
            UPDATE leave_accruals 
            SET hours_available = ?, last_updated = CURRENT_TIMESTAMP
            WHERE personnel = ? AND leave_type = ?
        """, (new_balance, personnel, leave_type))
    else:
        c.execute("""
            INSERT INTO leave_accruals (personnel, leave_type, hours_available, hours_used)
            VALUES (?, ?, ?, 0)
        """, (personnel, leave_type, hours))
    
    conn.commit()
    conn.close()
    return True

def set_accrual_balance(personnel, leave_type, exact_hours):
    """Set exact balance amount (not add/subtract, but SET to specific amount)"""
    conn = sqlite3.connect(DB_PATH)
    c = conn.cursor()
    
    # Check if accrual exists
    c.execute("SELECT id FROM leave_accruals WHERE personnel = ? AND leave_type = ?", 
              (personnel, leave_type))
    result = c.fetchone()
    
    if result:
        # Update to exact amount
        c.execute("""
            UPDATE leave_accruals 
            SET hours_available = ?, last_updated = CURRENT_TIMESTAMP
            WHERE personnel = ? AND leave_type = ?
        """, (exact_hours, personnel, leave_type))
    else:
        # Create new with exact amount
        c.execute("""
            INSERT INTO leave_accruals (personnel, leave_type, hours_available, hours_used)
            VALUES (?, ?, ?, 0)
        """, (personnel, leave_type, exact_hours))
    
    conn.commit()
    conn.close()
    return True

def deduct_leave_hours(personnel, leave_type, hours):
    """Deduct hours from leave balance (called when logging time off)"""
    conn = sqlite3.connect(DB_PATH)
    c = conn.cursor()
    
    # Get current balance
    c.execute("SELECT hours_available, hours_used FROM leave_accruals WHERE personnel = ? AND leave_type = ?", 
              (personnel, leave_type))
    result = c.fetchone()
    
    if result:
        current_available, current_used = result
        if current_available >= hours:
            new_available = current_available - hours
            new_used = current_used + hours
            c.execute("""
                UPDATE leave_accruals 
                SET hours_available = ?, hours_used = ?, last_updated = CURRENT_TIMESTAMP
                WHERE personnel = ? AND leave_type = ?
            """, (new_available, new_used, personnel, leave_type))
            conn.commit()
            conn.close()
            return True, new_available
        else:
            conn.close()
            return False, current_available
    else:
        # No accrual record - initialize with 0
        c.execute("""
            INSERT INTO leave_accruals (personnel, leave_type, hours_available, hours_used)
            VALUES (?, ?, 0, ?)
        """, (personnel, leave_type, hours))
        conn.commit()
        conn.close()
        return False, 0

def get_leave_balance_summary(personnel):
    """Get summary of all leave balances for a person"""
    conn = sqlite3.connect(DB_PATH)
    query = """
        SELECT la.leave_type, 
               la.hours_available, 
               la.hours_used,
               lt.default_annual_hours
        FROM leave_accruals la
        LEFT JOIN leave_types lt ON la.leave_type = lt.leave_type_name
        WHERE la.personnel = ?
        ORDER BY la.leave_type
    """
    df = pd.read_sql_query(query, conn, params=(personnel,))
    conn.close()
    return df

#############################################
# ROOM MANAGEMENT FUNCTIONS (NEW!)
#############################################

def get_room_numbers(active_only=True):
    """Get all room numbers"""
    conn = sqlite3.connect(DB_PATH)
    if active_only:
        df = pd.read_sql_query("SELECT * FROM room_numbers WHERE active = 1 ORDER BY room_number", conn)
    else:
        df = pd.read_sql_query("SELECT * FROM room_numbers ORDER BY active DESC, room_number", conn)
    conn.close()
    return df

def add_room_number(room_number, notes=None):
    """Add a new room"""
    conn = sqlite3.connect(DB_PATH)
    c = conn.cursor()
    try:
        c.execute("INSERT INTO room_numbers (room_number, notes) VALUES (?, ?)", (room_number, notes))
        conn.commit()
        conn.close()
        return True
    except sqlite3.IntegrityError:
        conn.close()
        return False

def update_room_number(room_id, **kwargs):
    """Update room details"""
    conn = sqlite3.connect(DB_PATH)
    c = conn.cursor()
    set_clause = ", ".join([f"{k} = ?" for k in kwargs.keys()])
    values = list(kwargs.values()) + [room_id]
    c.execute(f"UPDATE room_numbers SET {set_clause} WHERE id = ?", values)
    conn.commit()
    conn.close()

def toggle_room_active(room_id):
    """Toggle room active status"""
    conn = sqlite3.connect(DB_PATH)
    c = conn.cursor()
    c.execute("UPDATE room_numbers SET active = NOT active WHERE id = ?", (room_id,))
    conn.commit()
    conn.close()

# Equipment management functions
def add_equipment(name, serial_number, purchase_date, status, location, notes):
    conn = sqlite3.connect(DB_PATH)
    c = conn.cursor()
    c.execute('''
        INSERT INTO equipment (name, serial_number, purchase_date, status, location, notes)
        VALUES (?, ?, ?, ?, ?, ?)
    ''', (name, serial_number, purchase_date, status, location, notes))
    conn.commit()
    conn.close()

# Database operations
def add_activity(date, activity_type, hours, students_trained, personnel, equipment, course, room_number, time_start, time_end, turn_in, received, notes):
    conn = sqlite3.connect(DB_PATH)
    c = conn.cursor()
    c.execute('''
        INSERT INTO activities (date, activity_type, hours, students_trained, personnel, equipment, course, room_number, time_start, time_end, turn_in, received, notes)
        VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?)
    ''', (date, activity_type, hours, students_trained, personnel, equipment, course, room_number, time_start, time_end, turn_in, received, notes))
    conn.commit()
    conn.close()

def get_activities(start_date=None, end_date=None):
    conn = sqlite3.connect(DB_PATH)
    query = "SELECT * FROM activities"
    if start_date and end_date:
        query += f" WHERE date BETWEEN '{start_date}' AND '{end_date}'"
    query += " ORDER BY date DESC"
    df = pd.read_sql_query(query, conn)
    conn.close()
    return df

def delete_activity(activity_id):
    conn = sqlite3.connect(DB_PATH)
    c = conn.cursor()
    c.execute("DELETE FROM activities WHERE id = ?", (activity_id,))
    conn.commit()
    conn.close()

def update_activity(activity_id, date, activity_type, hours, students_trained, personnel, equipment, course, room_number, time_start, time_end, turn_in, received, notes):
    conn = sqlite3.connect(DB_PATH)
    c = conn.cursor()
    c.execute('''
        UPDATE activities 
        SET date=?, activity_type=?, hours=?, students_trained=?, personnel=?, equipment=?, course=?, room_number=?, time_start=?, time_end=?, turn_in=?, received=?, notes=?
        WHERE id=?
    ''', (date, activity_type, hours, students_trained, personnel, equipment, course, room_number, time_start, time_end, turn_in, received, notes, activity_id))
    conn.commit()
    conn.close()

# Personnel management
def add_personnel(name, role):
    conn = sqlite3.connect(DB_PATH)
    c = conn.cursor()
    try:
        c.execute("INSERT INTO personnel (name, role) VALUES (?, ?)", (name, role))
        conn.commit()
        success = True
    except sqlite3.IntegrityError:
        success = False
    conn.close()
    return success

def get_personnel(active_only=True):
    conn = sqlite3.connect(DB_PATH)
    query = "SELECT * FROM personnel"
    if active_only:
        query += " WHERE active = 1"
    query += " ORDER BY name"
    df = pd.read_sql_query(query, conn)
    conn.close()
    return df

def toggle_personnel(person_id, active):
    conn = sqlite3.connect(DB_PATH)
    c = conn.cursor()
    c.execute("UPDATE personnel SET active = ? WHERE id = ?", (active, person_id))
    conn.commit()
    conn.close()

# Equipment management
def add_equipment(name, status, notes):
    conn = sqlite3.connect(DB_PATH)
    c = conn.cursor()
    try:
        c.execute("INSERT INTO equipment (name, status, notes) VALUES (?, ?, ?)", (name, status, notes))
        conn.commit()
        success = True
    except sqlite3.IntegrityError:
        success = False
    conn.close()
    return success

def get_equipment(active_only=True):
    conn = sqlite3.connect(DB_PATH)
    query = "SELECT * FROM equipment"
    if active_only:
        query += " WHERE active = 1"
    query += " ORDER BY name"
    df = pd.read_sql_query(query, conn)
    conn.close()
    return df

def update_equipment_status(equipment_id, status, maintenance_date, notes):
    conn = sqlite3.connect(DB_PATH)
    c = conn.cursor()
    c.execute("""
        UPDATE equipment 
        SET status = ?, last_maintenance = ?, notes = ?
        WHERE id = ?
    """, (status, maintenance_date, notes, equipment_id))
    conn.commit()
    conn.close()

def toggle_equipment(equipment_id, active):
    conn = sqlite3.connect(DB_PATH)
    c = conn.cursor()
    c.execute("UPDATE equipment SET active = ? WHERE id = ?", (active, equipment_id))
    conn.commit()
    conn.close()

# Enhanced Course management with editing
def add_course(name, description):
    conn = sqlite3.connect(DB_PATH)
    c = conn.cursor()
    try:
        c.execute("""
            INSERT INTO courses (name, description) 
            VALUES (?, ?)
        """, (name, description))
        conn.commit()
        success = True
    except sqlite3.IntegrityError:
        success = False
    conn.close()
    return success

def update_course(course_id, name, description):
    conn = sqlite3.connect(DB_PATH)
    c = conn.cursor()
    c.execute("""
        UPDATE courses 
        SET name=?, description=?
        WHERE id=?
    """, (name, description, course_id))
    conn.commit()
    conn.close()

def get_courses(active_only=True):
    conn = sqlite3.connect(DB_PATH)
    query = "SELECT * FROM courses"
    if active_only:
        query += " WHERE active = 1"
    query += " ORDER BY name"
    df = pd.read_sql_query(query, conn)
    conn.close()
    return df

def toggle_course(course_id, active):
    conn = sqlite3.connect(DB_PATH)
    c = conn.cursor()
    c.execute("UPDATE courses SET active = ? WHERE id = ?", (active, course_id))
    conn.commit()
    conn.close()

# Enhanced Activity type management with editing
def add_activity_type(name, description):
    conn = sqlite3.connect(DB_PATH)
    c = conn.cursor()
    try:
        c.execute("INSERT INTO activity_types (name, description) VALUES (?, ?)", (name, description))
        conn.commit()
        success = True
    except sqlite3.IntegrityError:
        success = False
    conn.close()
    return success

def update_activity_type(type_id, name, description):
    conn = sqlite3.connect(DB_PATH)
    c = conn.cursor()
    c.execute("UPDATE activity_types SET name=?, description=? WHERE id=?", (name, description, type_id))
    conn.commit()
    conn.close()

def get_activity_types(active_only=True):
    conn = sqlite3.connect(DB_PATH)
    query = "SELECT * FROM activity_types"
    if active_only:
        query += " WHERE active = 1"
    query += " ORDER BY name"
    df = pd.read_sql_query(query, conn)
    conn.close()
    return df

def toggle_activity_type(type_id, active):
    conn = sqlite3.connect(DB_PATH)
    c = conn.cursor()
    c.execute("UPDATE activity_types SET active = ? WHERE id = ?", (active, type_id))
    conn.commit()
    conn.close()

# Incident management
def add_incident(date, incident_type, equipment, severity, description):
    conn = sqlite3.connect(DB_PATH)
    c = conn.cursor()
    c.execute("""
        INSERT INTO incidents (date, incident_type, equipment, severity, description)
        VALUES (?, ?, ?, ?, ?)
    """, (date, incident_type, equipment, severity, description))
    conn.commit()
    conn.close()

def get_incidents(resolved=None):
    conn = sqlite3.connect(DB_PATH)
    query = "SELECT * FROM incidents"
    if resolved is not None:
        query += f" WHERE resolved = {resolved}"
    query += " ORDER BY date DESC"
    df = pd.read_sql_query(query, conn)
    conn.close()
    return df

def resolve_incident(incident_id, resolution):
    conn = sqlite3.connect(DB_PATH)
    c = conn.cursor()
    c.execute("""
        UPDATE incidents 
        SET resolved = 1, resolution = ?
        WHERE id = ?
    """, (resolution, incident_id))
    conn.commit()
    conn.close()

# Goals management
def add_goal(goal_type, target_value, period):
    conn = sqlite3.connect(DB_PATH)
    c = conn.cursor()
    c.execute("""
        INSERT INTO goals (goal_type, target_value, period)
        VALUES (?, ?, ?)
    """, (goal_type, target_value, period))
    conn.commit()
    conn.close()

def get_goals():
    conn = sqlite3.connect(DB_PATH)
    df = pd.read_sql_query("SELECT * FROM goals ORDER BY created_at DESC", conn)
    conn.close()
    return df

def update_goal_progress(goal_id, current_value):
    conn = sqlite3.connect(DB_PATH)
    c = conn.cursor()
    c.execute("UPDATE goals SET current_value = ? WHERE id = ?", (current_value, goal_id))
    conn.commit()
    conn.close()

def delete_goal(goal_id):
    conn = sqlite3.connect(DB_PATH)
    c = conn.cursor()
    c.execute("DELETE FROM goals WHERE id = ?", (goal_id,))
    conn.commit()
    conn.close()

#############################################
# AI-POWERED FEATURES (PHASE 2)
#############################################

def generate_data_story(activities_df, period_name="selected period"):
    """Generate AI-powered narrative about the data using Claude API"""
    if activities_df.empty:
        return "No activities to analyze for this period."
    
    # Prepare data summary
    total_hours = activities_df['hours'].sum()
    total_activities = len(activities_df)
    total_students = activities_df['students_trained'].sum()
    
    # Activity type breakdown
    activity_type_summary = {}
    for idx, row in activities_df.iterrows():
        if pd.notna(row['activity_type']):
            types = [t.strip() for t in row['activity_type'].split(',')]
            for t in types:
                activity_type_summary[t] = activity_type_summary.get(t, 0) + row['hours']
    
    # Personnel breakdown
    personnel_hours = {}
    for idx, row in activities_df.iterrows():
        if pd.notna(row['personnel']):
            people = [p.strip() for p in row['personnel'].split(',')]
            for person in people:
                personnel_hours[person] = personnel_hours.get(person, 0) + row['hours']
    
    # Prepare prompt
    data_context = f"""
Data Summary for {period_name}:
- Total Hours: {total_hours:.1f}
- Total Activities: {total_activities}
- Students Trained: {int(total_students)}

Activity Type Distribution:
{json.dumps(activity_type_summary, indent=2)}

Top Personnel (by hours):
{json.dumps(dict(sorted(personnel_hours.items(), key=lambda x: x[1], reverse=True)[:5]), indent=2)}

Course Count: {activities_df['course'].notna().sum()} activities linked to courses
"""
    
    try:
        # Try to use Claude API if available (requires anthropic library and API key)
        import anthropic
        client = anthropic.Anthropic(api_key=st.secrets.get("ANTHROPIC_API_KEY", ""))
        
        message = client.messages.create(
            model="claude-sonnet-4-20250514",
            max_tokens=500,
            messages=[{
                "role": "user",
                "content": f"""You are analyzing simulation center operations data. Write a concise, insightful 3-4 sentence narrative summary highlighting the key patterns and notable points. Be specific with numbers. Focus on actionable insights.

{data_context}

Write in a professional but conversational tone suitable for leadership."""
            }]
        )
        
        return message.content[0].text
    
    except:
        # Fallback to rule-based narrative
        top_activity = max(activity_type_summary, key=activity_type_summary.get) if activity_type_summary else "Unknown"
        top_person = max(personnel_hours, key=personnel_hours.get) if personnel_hours else "Unknown"
        
        avg_duration = total_hours / total_activities if total_activities > 0 else 0
        
        story = f"During the {period_name}, your simulation center delivered {total_hours:.1f} hours across {total_activities} activities, training {int(total_students)} students. "
        story += f"{top_activity} was the primary activity type, accounting for {activity_type_summary.get(top_activity, 0):.1f} hours. "
        story += f"{top_person} led the team with {personnel_hours.get(top_person, 0):.1f} hours of contribution. "
        story += f"The average activity duration was {avg_duration:.2f} hours, indicating {'intensive' if avg_duration > 3 else 'standard'} session lengths."
        
        return story

def detect_anomalies(activities_df, historical_df=None):
    """Detect statistical anomalies and unusual patterns"""
    anomalies = []
    
    if activities_df.empty:
        return anomalies
    
    # Check for unusually long activities
    avg_hours = activities_df['hours'].mean()
    std_hours = activities_df['hours'].std()
    long_activities = activities_df[activities_df['hours'] > avg_hours + 2*std_hours]
    
    if not long_activities.empty:
        anomalies.append({
            'type': 'duration',
            'severity': 'info',
            'message': f"⏰ {len(long_activities)} activities were significantly longer than average ({long_activities['hours'].mean():.1f}h vs {avg_hours:.1f}h average)"
        })
    
    # Check for zero student count in training activities
    training_activities = activities_df[activities_df['activity_type'].str.contains('Training|Session', case=False, na=False)]
    zero_students = training_activities[training_activities['students_trained'] == 0]
    
    if not zero_students.empty:
        anomalies.append({
            'type': 'students',
            'severity': 'warning',
            'message': f"👥 {len(zero_students)} training activities recorded with 0 students - verify data accuracy"
        })
    
    # Check for unusual activity patterns by day of week
    if len(activities_df) >= 7:
        activities_df['day_of_week'] = pd.to_datetime(activities_df['date']).dt.dayofweek
        day_distribution = activities_df.groupby('day_of_week').size()
        
        if day_distribution.std() > day_distribution.mean():
            busiest_day = ['Monday', 'Tuesday', 'Wednesday', 'Thursday', 'Friday', 'Saturday', 'Sunday'][day_distribution.idxmax()]
            anomalies.append({
                'type': 'pattern',
                'severity': 'info',
                'message': f"📅 Activity distribution is uneven - {busiest_day} has significantly more activities than other days"
            })
    
    # Check for equipment not used
    if historical_df is not None and not historical_df.empty:
        recent_equipment = set()
        for idx, row in activities_df.iterrows():
            if pd.notna(row['equipment']):
                recent_equipment.update([e.strip() for e in row['equipment'].split(',')])
        
        historical_equipment = set()
        for idx, row in historical_df.iterrows():
            if pd.notna(row['equipment']):
                historical_equipment.update([e.strip() for e in row['equipment'].split(',')])
        
        unused = historical_equipment - recent_equipment
        if unused and len(unused) <= 5:
            anomalies.append({
                'type': 'equipment',
                'severity': 'warning',
                'message': f"🔧 Equipment not used in current period: {', '.join(list(unused)[:3])}"
            })
    
    return anomalies

def generate_actionable_insights(activities_df):
    """Generate specific actionable recommendations based on data patterns"""
    insights = []
    
    if activities_df.empty:
        return insights
    
    # Room utilization insights
    if activities_df['room_number'].notna().any():
        room_usage = {}
        for idx, row in activities_df.iterrows():
            if pd.notna(row['room_number']):
                rooms = [r.strip() for r in row['room_number'].split(',')]
                for room in rooms:
                    room_usage[room] = room_usage.get(room, 0) + 1
        
        if room_usage:
            avg_usage = sum(room_usage.values()) / len(room_usage)
            underused_rooms = [room for room, count in room_usage.items() if count < avg_usage * 0.5]
            
            if underused_rooms:
                insights.append({
                    'category': 'Space Optimization',
                    'icon': '🏢',
                    'recommendation': f"Rooms {', '.join(underused_rooms[:2])} are underutilized. Consider consolidating activities or promoting these spaces for scheduling.",
                    'impact': 'medium'
                })
    
    # Scheduling optimization
    if len(activities_df) >= 5:
        activities_df['hour'] = pd.to_datetime(activities_df['time_start'], format='%H:%M', errors='coerce').dt.hour
        hour_counts = activities_df['hour'].value_counts()
        
        if not hour_counts.empty:
            peak_hour = hour_counts.idxmax()
            if pd.notna(peak_hour):
                insights.append({
                    'category': 'Scheduling',
                    'icon': '⏰',
                    'recommendation': f"Peak activity time is around {int(peak_hour)}:00. Consider staffing adjustments during this window for optimal support.",
                    'impact': 'high'
                })
    
    # Course efficiency
    course_stats = activities_df[activities_df['course'].notna()].groupby('course').agg({
        'hours': 'mean',
        'students_trained': 'mean'
    })
    
    if not course_stats.empty:
        low_efficiency = course_stats[course_stats['students_trained'] < course_stats['students_trained'].mean() * 0.7]
        if not low_efficiency.empty:
            insights.append({
                'category': 'Course Efficiency',
                'icon': '📚',
                'recommendation': f"{len(low_efficiency)} courses have below-average student counts. Review {low_efficiency.index[0]} for potential scheduling or marketing improvements.",
                'impact': 'medium'
            })
    
    # Personnel workload balance
    personnel_hours = {}
    for idx, row in activities_df.iterrows():
        if pd.notna(row['personnel']):
            people = [p.strip() for p in row['personnel'].split(',')]
            for person in people:
                personnel_hours[person] = personnel_hours.get(person, 0) + row['hours']
    
    if len(personnel_hours) >= 3:
        hours_list = list(personnel_hours.values())
        max_hours = max(hours_list)
        min_hours = min(hours_list)
        
        if max_hours > min_hours * 2:
            heaviest = max(personnel_hours, key=personnel_hours.get)
            insights.append({
                'category': 'Workload Balance',
                'icon': '⚖️',
                'recommendation': f"Workload imbalance detected. {heaviest} has 2x+ hours of others. Consider redistributing responsibilities to prevent burnout.",
                'impact': 'high'
            })
    
    # Equipment maintenance reminder
    equipment_usage = {}
    for idx, row in activities_df.iterrows():
        if pd.notna(row['equipment']):
            items = [e.strip() for e in row['equipment'].split(',')]
            for item in items:
                equipment_usage[item] = equipment_usage.get(item, 0) + 1
    
    if equipment_usage:
        heavily_used = [eq for eq, count in equipment_usage.items() if count > len(activities_df) * 0.3]
        if heavily_used:
            insights.append({
                'category': 'Equipment Maintenance',
                'icon': '🔧',
                'recommendation': f"{heavily_used[0]} is heavily utilized ({equipment_usage[heavily_used[0]]} uses). Schedule preventive maintenance to avoid downtime.",
                'impact': 'high'
            })
    
    return insights

def create_activity_heatmap(activities_df):
    """Create heatmap showing activity intensity by day and hour"""
    if activities_df.empty:
        return None
    
    # Parse datetime
    activities_df = activities_df.copy()
    activities_df['datetime'] = pd.to_datetime(activities_df['date'])
    activities_df['day_of_week'] = activities_df['datetime'].dt.day_name()
    
    # Parse hour from time_start
    activities_df['hour'] = pd.to_datetime(activities_df['time_start'], format='%H:%M', errors='coerce').dt.hour
    
    # Filter valid hours
    valid_data = activities_df[activities_df['hour'].notna()].copy()
    
    if valid_data.empty:
        return None
    
    # Create heatmap data
    heatmap_data = valid_data.groupby(['day_of_week', 'hour']).size().reset_index(name='count')
    
    # Pivot for heatmap
    days_order = ['Monday', 'Tuesday', 'Wednesday', 'Thursday', 'Friday', 'Saturday', 'Sunday']
    heatmap_pivot = heatmap_data.pivot(index='day_of_week', columns='hour', values='count').fillna(0)
    heatmap_pivot = heatmap_pivot.reindex(days_order)
    
    # Create heatmap
    fig = go.Figure(data=go.Heatmap(
        z=heatmap_pivot.values,
        x=heatmap_pivot.columns,
        y=heatmap_pivot.index,
        colorscale='YlOrRd',
        text=heatmap_pivot.values,
        texttemplate='%{text}',
        textfont={"size": 10},
        colorbar=dict(title="Activities")
    ))
    
    fig.update_layout(
        title="Activity Intensity Heatmap",
        xaxis_title="Hour of Day",
        yaxis_title="Day of Week",
        height=400
    )
    
    return fig

def create_sankey_diagram(activities_df):
    """Create Sankey diagram showing flow from activity type → personnel → equipment"""
    if activities_df.empty:
        return None
    
    # Build relationships
    relationships = []
    
    for idx, row in activities_df.iterrows():
        if pd.notna(row['activity_type']):
            activity_types = [t.strip() for t in row['activity_type'].split(',')]
            personnel_list = [p.strip() for p in row['personnel'].split(',')] if pd.notna(row['personnel']) else []
            equipment_list = [e.strip() for e in row['equipment'].split(',')] if pd.notna(row['equipment']) else []
            
            for act_type in activity_types:
                for person in personnel_list:
                    relationships.append(('Activity', act_type, person, row['hours']))
                    
                    for equipment in equipment_list:
                        relationships.append(('Personnel', person, equipment, row['hours']))
    
    if not relationships:
        return None
    
    # Aggregate relationships
    from collections import defaultdict
    aggregated = defaultdict(float)
    
    for layer, source, target, value in relationships:
        aggregated[(layer, source, target)] += value
    
    # Create nodes
    all_nodes = set()
    for (layer, source, target), value in aggregated.items():
        all_nodes.add(source)
        all_nodes.add(target)
    
    node_list = list(all_nodes)
    node_dict = {node: idx for idx, node in enumerate(node_list)}
    
    # Create links
    sources = []
    targets = []
    values = []
    
    for (layer, source, target), value in aggregated.items():
        sources.append(node_dict[source])
        targets.append(node_dict[target])
        values.append(value)
    
    # Limit to top flows
    if len(sources) > 50:
        top_indices = sorted(range(len(values)), key=lambda i: values[i], reverse=True)[:50]
        sources = [sources[i] for i in top_indices]
        targets = [targets[i] for i in top_indices]
        values = [values[i] for i in top_indices]
    
    # Create Sankey
    fig = go.Figure(data=[go.Sankey(
        node=dict(
            pad=15,
            thickness=20,
            line=dict(color="black", width=0.5),
            label=node_list,
            color='#F8B400'
        ),
        link=dict(
            source=sources,
            target=targets,
            value=values,
            color='rgba(248, 180, 0, 0.3)'
        )
    )])
    
    fig.update_layout(
        title="Activity Flow: Type → Personnel → Equipment",
        height=500,
        font=dict(size=10)
    )
    
    return fig

def create_powerpoint_report(activities_df, period_name, start_date, end_date):
    """Create PowerPoint presentation with executive briefing"""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
        from pptx.dml.color import RGBColor
    except ImportError:
        return None, "python-pptx library not installed. Install with: pip install python-pptx --break-system-packages"
    
    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # VCU Colors
    vcu_gold = RGBColor(248, 180, 0)
    vcu_black = RGBColor(0, 0, 0)
    
    # SLIDE 1: Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "VCU Health Sciences\nSimulation Center Operations Report"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = vcu_gold
    title_para.alignment = PP_ALIGN.CENTER
    
    # Period
    period_box = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(8), Inches(0.5))
    period_frame = period_box.text_frame
    period_frame.text = f"Period: {period_name}"
    period_para = period_frame.paragraphs[0]
    period_para.font.size = Pt(24)
    period_para.alignment = PP_ALIGN.CENTER
    
    # Date generated
    date_box = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(8), Inches(0.5))
    date_frame = date_box.text_frame
    date_frame.text = f"Generated: {datetime.now().strftime('%B %d, %Y')}"
    date_para = date_frame.paragraphs[0]
    date_para.font.size = Pt(14)
    date_para.alignment = PP_ALIGN.CENTER
    
    # SLIDE 2: Executive Summary
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Title only
    title = slide.shapes.title
    title.text = "Executive Summary"
    title.text_frame.paragraphs[0].font.color.rgb = vcu_gold
    
    # Calculate metrics
    total_hours = activities_df['hours'].sum()
    total_activities = len(activities_df)
    total_students = activities_df['students_trained'].sum()
    unique_courses = activities_df['course'].dropna().nunique()
    
    # Add metrics text box
    left = Inches(1)
    top = Inches(2)
    width = Inches(8)
    height = Inches(4)
    
    text_box = slide.shapes.add_textbox(left, top, width, height)
    tf = text_box.text_frame
    tf.word_wrap = True
    
    # Add metrics
    metrics_text = f"""
Total Operating Hours: {total_hours:.0f} hours
Total Activities: {total_activities}
Students Trained: {int(total_students):,}
Courses Delivered: {unique_courses}
Average Session Duration: {activities_df['hours'].mean():.1f} hours
Average Class Size: {(total_students / total_activities):.1f} students
"""
    
    tf.text = metrics_text.strip()
    for paragraph in tf.paragraphs:
        paragraph.font.size = Pt(20)
        paragraph.space_after = Pt(12)
    
    # SLIDE 3: AI Narrative
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    title.text = "Key Insights"
    title.text_frame.paragraphs[0].font.color.rgb = vcu_gold
    
    # Generate story
    story = generate_data_story(activities_df, period_name)
    
    text_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4.5))
    tf = text_box.text_frame
    tf.word_wrap = True
    tf.text = story
    
    for paragraph in tf.paragraphs:
        paragraph.font.size = Pt(18)
        paragraph.space_after = Pt(12)
        paragraph.alignment = PP_ALIGN.LEFT
    
    # SLIDE 4: Activity Type Breakdown
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    title.text = "Activity Type Distribution"
    title.text_frame.paragraphs[0].font.color.rgb = vcu_gold
    
    # Parse activity types
    activity_type_hours = {}
    for idx, row in activities_df.iterrows():
        if pd.notna(row['activity_type']):
            types = [t.strip() for t in row['activity_type'].split(',')]
            for act_type in types:
                activity_type_hours[act_type] = activity_type_hours.get(act_type, 0) + row['hours']
    
    # Create table
    rows = len(activity_type_hours) + 1
    cols = 3
    left = Inches(2)
    top = Inches(2)
    width = Inches(6)
    height = Inches(0.5) * rows
    
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    
    # Header
    table.cell(0, 0).text = "Activity Type"
    table.cell(0, 1).text = "Hours"
    table.cell(0, 2).text = "Percentage"
    
    for cell in [table.cell(0, i) for i in range(cols)]:
        cell.fill.solid()
        cell.fill.fore_color.rgb = vcu_gold
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    
    # Data
    total_activity_hours = sum(activity_type_hours.values())
    for idx, (act_type, hours) in enumerate(sorted(activity_type_hours.items(), key=lambda x: x[1], reverse=True)):
        table.cell(idx + 1, 0).text = act_type
        table.cell(idx + 1, 1).text = f"{hours:.1f}"
        table.cell(idx + 1, 2).text = f"{(hours/total_activity_hours*100):.1f}%"
    
    # SLIDE 5: Actionable Insights
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    title.text = "Actionable Recommendations"
    title.text_frame.paragraphs[0].font.color.rgb = vcu_gold
    
    insights = generate_actionable_insights(activities_df)
    
    text_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4.5))
    tf = text_box.text_frame
    tf.word_wrap = True
    
    if insights:
        for insight in insights[:5]:  # Top 5 insights
            p = tf.add_paragraph()
            p.text = f"{insight['icon']} {insight['category']}: {insight['recommendation']}"
            p.font.size = Pt(16)
            p.space_after = Pt(12)
            p.level = 0
    else:
        tf.text = "✅ No critical action items. Operations performing optimally."
        tf.paragraphs[0].font.size = Pt(18)
    
    # SLIDE 6: Top Courses
    if not activities_df[activities_df['course'].notna()].empty:
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        title = slide.shapes.title
        title.text = "Top Performing Courses"
        title.text_frame.paragraphs[0].font.color.rgb = vcu_gold
        
        top_courses = activities_df[activities_df['course'].notna()].groupby('course').agg({
            'hours': 'sum',
            'students_trained': 'sum'
        }).sort_values('hours', ascending=False).head(5)
        
        rows = len(top_courses) + 1
        cols = 3
        
        table = slide.shapes.add_table(rows, cols, Inches(2), Inches(2), Inches(6), Inches(0.5) * rows).table
        
        # Header
        table.cell(0, 0).text = "Course"
        table.cell(0, 1).text = "Total Hours"
        table.cell(0, 2).text = "Students"
        
        for cell in [table.cell(0, i) for i in range(cols)]:
            cell.fill.solid()
            cell.fill.fore_color.rgb = vcu_gold
            cell.text_frame.paragraphs[0].font.bold = True
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        
        # Data
        for idx, (course, row) in enumerate(top_courses.iterrows()):
            table.cell(idx + 1, 0).text = course
            table.cell(idx + 1, 1).text = f"{row['hours']:.1f}"
            table.cell(idx + 1, 2).text = str(int(row['students_trained']))
    
    # SLIDE 7: Next Steps
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    title.text = "Next Steps"
    title.text_frame.paragraphs[0].font.color.rgb = vcu_gold
    
    text_box = slide.shapes.add_textbox(Inches(1.5), Inches(2.5), Inches(7), Inches(3))
    tf = text_box.text_frame
    tf.word_wrap = True
    
    next_steps = """
• Review and address high-priority recommendations
• Continue monitoring key performance indicators
• Maintain focus on quality training delivery
• Schedule follow-up analysis for next period
"""
    
    tf.text = next_steps.strip()
    for paragraph in tf.paragraphs:
        paragraph.font.size = Pt(20)
        paragraph.space_after = Pt(16)
    
    # Save to BytesIO
    pptx_io = BytesIO()
    prs.save(pptx_io)
    pptx_io.seek(0)
    
    return pptx_io, None

#############################################
# PDF REPORT GENERATION FUNCTIONS
#############################################

# Initialize database
init_db()

# Custom CSS for VCU branding and mobile responsiveness
st.markdown("""
    <style>
    .main {
        background-color: #f5f5f5;
    }
    .stMetric {
        background-color: white;
        padding: 15px;
        border-radius: 5px;
        box-shadow: 0 2px 4px rgba(0,0,0,0.1);
    }
    .metric-card {
        background: linear-gradient(135deg, #F8B400 0%, #000000 100%);
        padding: 20px;
        border-radius: 10px;
        color: white;
        margin: 10px 0;
    }
    
    /* Mobile Responsiveness */
    @media (max-width: 768px) {
        .main .block-container {
            padding-left: 1rem;
            padding-right: 1rem;
        }
        .stButton button {
            width: 100%;
            margin: 5px 0;
        }
        .stSelectbox, .stMultiSelect {
            margin-bottom: 1rem;
        }
        .row-widget {
            flex-direction: column;
        }
    }
    
    /* Enhanced Button Styling */
    .stButton > button {
        border-radius: 8px;
        font-weight: 500;
        transition: all 0.3s ease;
    }
    
    .stButton > button:hover {
        transform: translateY(-2px);
        box-shadow: 0 4px 8px rgba(0,0,0,0.2);
    }
    
    /* Improved Card Styling */
    div[data-testid="stExpander"] {
        background-color: white;
        border-radius: 8px;
        border: 1px solid #e0e0e0;
        margin: 10px 0;
    }
    
    /* Better Table Styling */
    .dataframe {
        border-radius: 8px;
        overflow: hidden;
    }
    
    /* Tooltip Enhancement */
    [data-baseweb="tooltip"] {
        background-color: #333;
        color: white;
        padding: 8px 12px;
        border-radius: 6px;
        font-size: 14px;
    }
    
    /* Improved Chart Containers */
    .js-plotly-plot {
        border-radius: 8px;
        background: white;
        padding: 10px;
        box-shadow: 0 2px 4px rgba(0,0,0,0.05);
    }
    </style>
""", unsafe_allow_html=True)

# Sidebar navigation
with st.sidebar:
    st.image("https://www.vcuhealth.org/sites/default/files/VCU-Health-Logo.svg", width=200)
    st.markdown("---")
    
    page = st.radio(
        "Navigation",
        ["📝 Data Entry", "📊 Dashboard", "👔 Executive Dashboard", "📚 Course Analytics", "🔧 Equipment Analytics", 
         "📅 History", "❌ Cancellations", "🏖️ Time Off", "🛠️ Equipment", "⚠️ Incidents", "🎯 Goals", "⚙️ Settings"],
        key="nav"
    )
    
    st.markdown("---")
    st.markdown("### Quick Stats")
    
    # Time period selector
    time_period = st.selectbox(
        "Period",
        ["Today", "This Week", "This Month", "This Quarter", "This Year"],
        index=1,  # Default to This Week
        key="quick_stats_period"
    )
    
    # Calculate date ranges based on selection
    today = datetime.now().date()
    
    if time_period == "Today":
        start_date = today
        end_date = today
        prev_start = today - timedelta(days=1)
        prev_end = today - timedelta(days=1)
        comparison_label = "vs. Yesterday"
    elif time_period == "This Week":
        start_date = today - timedelta(days=today.weekday())
        end_date = start_date + timedelta(days=6)
        prev_start = start_date - timedelta(days=7)
        prev_end = end_date - timedelta(days=7)
        comparison_label = "vs. Last Week"
    elif time_period == "This Month":
        start_date = today.replace(day=1)
        # Last day of current month
        if today.month == 12:
            end_date = today.replace(day=31)
        else:
            end_date = (today.replace(month=today.month+1, day=1) - timedelta(days=1))
        # Previous month
        if today.month == 1:
            prev_start = today.replace(year=today.year-1, month=12, day=1)
            prev_end = today.replace(year=today.year-1, month=12, day=31)
        else:
            prev_start = today.replace(month=today.month-1, day=1)
            prev_end = start_date - timedelta(days=1)
        comparison_label = "vs. Last Month"
    elif time_period == "This Quarter":
        quarter = (today.month - 1) // 3
        start_date = today.replace(month=quarter*3+1, day=1)
        end_month = quarter*3+3
        if end_month == 12:
            end_date = today.replace(month=12, day=31)
        else:
            end_date = today.replace(month=end_month+1, day=1) - timedelta(days=1)
        # Previous quarter
        prev_quarter = quarter - 1 if quarter > 0 else 3
        prev_year = today.year if quarter > 0 else today.year - 1
        prev_start = today.replace(year=prev_year, month=prev_quarter*3+1, day=1)
        prev_end_month = prev_quarter*3+3
        if prev_end_month == 12:
            prev_end = today.replace(year=prev_year, month=12, day=31)
        else:
            prev_end = today.replace(year=prev_year, month=prev_end_month+1, day=1) - timedelta(days=1)
        comparison_label = "vs. Last Quarter"
    else:  # This Year
        start_date = today.replace(month=1, day=1)
        end_date = today.replace(month=12, day=31)
        prev_start = today.replace(year=today.year-1, month=1, day=1)
        prev_end = today.replace(year=today.year-1, month=12, day=31)
        comparison_label = "vs. Last Year"
    
    # Get current period activities (EXCLUDING CANCELLED)
    current_activities = get_activities(start_date, end_date)
    if not current_activities.empty:
        current_activities = current_activities[~current_activities['notes'].str.contains(r'\[CANCELLED:', na=False, regex=True)]
    
    # Get previous period activities for comparison (EXCLUDING CANCELLED)
    prev_activities = get_activities(prev_start, prev_end)
    if not prev_activities.empty:
        prev_activities = prev_activities[~prev_activities['notes'].str.contains(r'\[CANCELLED:', na=False, regex=True)]
    
    if not current_activities.empty:
        total_hours = current_activities['hours'].sum()
        total_activities = len(current_activities)
        
        # Calculate changes
        prev_hours = prev_activities['hours'].sum() if not prev_activities.empty else 0
        prev_count = len(prev_activities) if not prev_activities.empty else 0
        
        hours_delta = total_hours - prev_hours
        activities_delta = total_activities - prev_count
        
        st.metric(
            f"{time_period} - Hours", 
            f"{total_hours:.1f}",
            f"{hours_delta:+.1f} {comparison_label}",
            delta_color="normal"
        )
        st.metric(
            f"{time_period} - Activities", 
            total_activities,
            f"{activities_delta:+d} {comparison_label}",
            delta_color="normal"
        )
        
        # Tech Hours - sum hours for specific tech team members
        tech_team = ["Hayden", "Justin", "Freddie", "Leana", "Tony"]
        tech_hours = 0
        for idx, row in current_activities.iterrows():
            if pd.notna(row['personnel']) and row['personnel']:
                personnel_list = [p.strip() for p in row['personnel'].split(',')]
                if any(tech in personnel_list for tech in tech_team):
                    tech_hours += row['hours']
        
        # Previous period tech hours
        prev_tech_hours = 0
        for idx, row in prev_activities.iterrows():
            if pd.notna(row['personnel']) and row['personnel']:
                personnel_list = [p.strip() for p in row['personnel'].split(',')]
                if any(tech in personnel_list for tech in tech_team):
                    prev_tech_hours += row['hours']
        
        tech_delta = tech_hours - prev_tech_hours
        
        st.metric(
            f"🔧 Tech Hours",
            f"{tech_hours:.1f}",
            f"{tech_delta:+.1f} {comparison_label}",
            delta_color="normal"
        )
    else:
        st.metric(f"{time_period} - Hours", "0", "No previous data")
        st.metric(f"{time_period} - Activities", "0", "No previous data")
        st.metric("🔧 Tech Hours", "0", "No previous data")

#############################################
# PAGE 1: DATA ENTRY
#############################################
if page == "📝 Data Entry":
    st.title("📝 Activity Data Entry")
    
    col1, col2 = st.columns(2)
    
    with col1:
        entry_date = st.date_input("Date *", value=datetime.now())
        
        activity_types = get_activity_types()
        selected_activity_types = st.multiselect(
            "Activity Type(s) *",
            activity_types['name'].tolist() if not activity_types.empty else []
        )
        
        hours_worked = st.number_input("Hours *", min_value=0.0, max_value=24.0, step=0.25, value=1.0)
        students_trained = st.number_input("Students Trained", min_value=0, step=1, value=0)
    
    with col2:
        personnel_df = get_personnel()
        personnel_list = personnel_df['name'].tolist() if not personnel_df.empty else []
        selected_personnel = st.multiselect("Personnel", personnel_list)
        
        courses_df = get_courses()
        course_list = [""] + courses_df['name'].tolist() if not courses_df.empty else [""]
        selected_course = st.selectbox("Course", course_list)
        
        equipment_df = get_equipment()
        equipment_list = equipment_df['name'].tolist() if not equipment_df.empty else []
        selected_equipment = st.multiselect("Equipment Used", equipment_list)
    
    # NEW: Room and Time fields
    st.markdown("### 📍 Location & Schedule")
    col1, col2, col3 = st.columns(3)
    
    with col1:
        selected_rooms = st.multiselect("Room Number(s)", get_active_rooms())
    
    with col2:
        time_start = st.time_input("Start Time", value=None)
    
    with col3:
        time_end = st.time_input("End Time", value=None)
    
    st.markdown("---")
    
    col1, col2 = st.columns(2)
    with col1:
        turn_in = st.number_input("Turn In Count", min_value=0, step=1, value=0)
    with col2:
        received = st.number_input("Received Count", min_value=0, step=1, value=0)
    
    notes = st.text_area("Notes", height=100)
    
    if st.button("💾 Save Activity", type="primary", use_container_width=True):
        if entry_date and selected_activity_types and hours_worked > 0:
            activity_type_str = ", ".join(selected_activity_types)
            personnel_str = ", ".join(selected_personnel) if selected_personnel else ""
            equipment_str = ", ".join(selected_equipment) if selected_equipment else ""
            room_str = ", ".join(selected_rooms) if selected_rooms else ""
            time_start_str = time_start.strftime("%H:%M") if time_start else ""
            time_end_str = time_end.strftime("%H:%M") if time_end else ""
            
            add_activity(
                entry_date, activity_type_str, hours_worked, students_trained,
                personnel_str, equipment_str, selected_course, room_str,
                time_start_str, time_end_str, turn_in, received, notes
            )
            st.success("✅ Activity saved successfully!")
            st.rerun()
        else:
            st.error("Please fill in all required fields (*)")

#############################################
# PAGE 2: DASHBOARD
#############################################
elif page == "📊 Dashboard":
    st.title("📊 Performance Dashboard")
    
    # Date range selector
    col1, col2, col3 = st.columns(3)
    with col1:
        date_range = st.selectbox(
            "View Period",
            ["This Week", "Last Week", "This Month", "Last Month", "Last 30 Days", "Custom"]
        )
    
    today = datetime.now().date()
    
    if date_range == "This Week":
        start_date = today - timedelta(days=today.weekday())
        end_date = today
    elif date_range == "Last Week":
        start_date = today - timedelta(days=today.weekday() + 7)
        end_date = start_date + timedelta(days=6)
    elif date_range == "This Month":
        start_date = today.replace(day=1)
        end_date = today
    elif date_range == "Last Month":
        first_this_month = today.replace(day=1)
        end_date = first_this_month - timedelta(days=1)
        start_date = end_date.replace(day=1)
    elif date_range == "Last 30 Days":
        start_date = today - timedelta(days=30)
        end_date = today
    else:  # Custom
        with col2:
            start_date = st.date_input("Start Date", value=today - timedelta(days=30))
        with col3:
            end_date = st.date_input("End Date", value=today)
    
    activities = get_activities(start_date, end_date)
    
    if not activities.empty:
        # Key Metrics
        col1, col2, col3, col4 = st.columns(4)
        
        with col1:
            total_hours = activities['hours'].sum()
            st.metric("⏰ Total Hours", f"{total_hours:.1f}")
        
        with col2:
            total_activities = len(activities)
            st.metric("📋 Activities", total_activities)
        
        with col3:
            unique_courses = activities['course'].dropna().nunique()
            st.metric("📚 Courses Run", unique_courses)
        
        with col4:
            avg_hours = activities['hours'].mean()
            st.metric("⌛ Avg Hours/Activity", f"{avg_hours:.2f}")
        
        st.markdown("---")
        
        # AI-POWERED DATA STORY (PHASE 2 FEATURE 1)
        st.markdown("### 📖 AI-Generated Summary")
        with st.spinner("Generating intelligent narrative..."):
            story = generate_data_story(activities, period_name=date_range.lower())
            st.info(story)
        
        # ACTIONABLE INSIGHTS (PHASE 2 FEATURE 2)
        st.markdown("---")
        st.markdown("### 💡 Actionable Recommendations")
        
        insights = generate_actionable_insights(activities)
        
        if insights:
            for insight in insights:
                if insight['impact'] == 'high':
                    st.error(f"{insight['icon']} **{insight['category']}**: {insight['recommendation']}")
                elif insight['impact'] == 'medium':
                    st.warning(f"{insight['icon']} **{insight['category']}**: {insight['recommendation']}")
                else:
                    st.info(f"{insight['icon']} **{insight['category']}**: {insight['recommendation']}")
        else:
            st.success("✅ No critical recommendations at this time. Operations are running smoothly!")
        
        # ANOMALY DETECTION (PHASE 2 FEATURE 4)
        st.markdown("---")
        st.markdown("### 🔍 Anomaly Detection")
        
        # Get historical data for comparison (past 90 days)
        historical_start = start_date - timedelta(days=90)
        historical_activities = get_activities(historical_start, start_date)
        
        anomalies = detect_anomalies(activities, historical_activities)
        
        if anomalies:
            for anomaly in anomalies:
                if anomaly['severity'] == 'warning':
                    st.warning(f"⚠️ {anomaly['message']}")
                else:
                    st.info(f"ℹ️ {anomaly['message']}")
        else:
            st.success("✅ No anomalies detected. All metrics within normal ranges.")
        
        st.markdown("---")
        
        # Charts
        col1, col2 = st.columns(2)
        
        with col1:
            st.markdown("### 📊 Hours by Activity Type")
            
            # Parse activity types from comma-separated strings
            activity_type_hours = {}
            for idx, row in activities.iterrows():
                if pd.notna(row['activity_type']) and row['activity_type']:
                    types = [t.strip() for t in row['activity_type'].split(',')]
                    for activity_type in types:
                        if activity_type:
                            activity_type_hours[activity_type] = activity_type_hours.get(activity_type, 0) + row['hours']
            
            if activity_type_hours:
                activity_summary = pd.Series(activity_type_hours).sort_values(ascending=False)
                
                fig = px.bar(
                    x=activity_summary.values,
                    y=activity_summary.index,
                    orientation='h',
                    labels={'x': 'Hours', 'y': 'Activity Type'},
                    color=activity_summary.index,
                    color_discrete_sequence=['#4A90E2', '#27AE60', '#9B59B6', '#E67E22', '#E74C3C', '#17A2B8', '#F8B400', '#34495E']
                )
                fig.update_layout(showlegend=False, height=400)
                st.plotly_chart(fig, use_container_width=True)
        
        with col2:
            st.markdown("### 🥧 Activity Distribution")
            
            if activity_type_hours:
                activity_summary = pd.Series(activity_type_hours)
                
                fig = px.pie(
                    values=activity_summary.values,
                    names=activity_summary.index,
                    color_discrete_sequence=['#4A90E2', '#27AE60', '#9B59B6', '#E67E22', '#E74C3C', '#17A2B8', '#F8B400', '#34495E']
                )
                fig.update_traces(textposition='inside', textinfo='percent+label')
                fig.update_layout(height=400, showlegend=True, legend=dict(orientation="v", yanchor="middle", y=0.5, xanchor="left", x=1.02))
                st.plotly_chart(fig, use_container_width=True)
        
        # Interactive Drill-Down Section
        st.markdown("---")
        st.markdown("### 🔍 Interactive Drill-Down")
        
        col1, col2, col3 = st.columns(3)
        
        with col1:
            # Activity Type Filter
            all_activity_types = ["All"] + sorted(list(activity_type_hours.keys())) if activity_type_hours else ["All"]
            selected_activity_type = st.selectbox("Filter by Activity Type", all_activity_types, key="drill_activity")
        
        with col2:
            # Course Filter
            course_list = ["All"] + sorted(activities['course'].dropna().unique().tolist())
            selected_course = st.selectbox("Filter by Course", course_list, key="drill_course")
        
        with col3:
            # Personnel Filter
            all_personnel = set()
            for idx, row in activities.iterrows():
                if pd.notna(row['personnel']) and row['personnel']:
                    all_personnel.update([p.strip() for p in row['personnel'].split(',')])
            personnel_list = ["All"] + sorted(list(all_personnel))
            selected_personnel = st.selectbox("Filter by Personnel", personnel_list, key="drill_personnel")
        
        # Apply filters
        filtered_activities = activities.copy()
        
        if selected_activity_type != "All":
            filtered_activities = filtered_activities[
                filtered_activities['activity_type'].apply(
                    lambda x: selected_activity_type in [t.strip() for t in str(x).split(',')] if pd.notna(x) else False
                )
            ]
        
        if selected_course != "All":
            filtered_activities = filtered_activities[filtered_activities['course'] == selected_course]
        
        if selected_personnel != "All":
            filtered_activities = filtered_activities[
                filtered_activities['personnel'].apply(
                    lambda x: selected_personnel in [p.strip() for p in str(x).split(',')] if pd.notna(x) else False
                )
            ]
        
        # Display filtered results
        if not filtered_activities.empty:
            st.markdown(f"**Showing {len(filtered_activities)} activities** ({filtered_activities['hours'].sum():.1f} hours)")
            
            # Show summary stats for filtered data
            col1, col2, col3, col4 = st.columns(4)
            with col1:
                st.metric("Hours", f"{filtered_activities['hours'].sum():.1f}")
            with col2:
                st.metric("Activities", len(filtered_activities))
            with col3:
                st.metric("Students", int(filtered_activities['students_trained'].sum()))
            with col4:
                avg_hours = filtered_activities['hours'].mean()
                st.metric("Avg Hours", f"{avg_hours:.2f}")
            
            # Show detailed table with expander
            with st.expander("📋 View Detailed Activities", expanded=False):
                display_df = filtered_activities[['date', 'activity_type', 'course', 'hours', 'students_trained', 'personnel', 'room_number']].copy()
                display_df['date'] = pd.to_datetime(display_df['date']).dt.strftime('%Y-%m-%d')
                st.dataframe(display_df, use_container_width=True, height=400)
        else:
            st.info("No activities match the selected filters.")
        
        # Daily Trend
        st.markdown("---")
        st.markdown("### 📈 Daily Activity Trend")
        daily_hours = activities.groupby('date')['hours'].sum().reset_index()
        
        fig = px.line(
            daily_hours,
            x='date',
            y='hours',
            markers=True,
            line_shape='spline'
        )
        fig.update_traces(line_color='#F8B400', marker=dict(color='#000000', size=8))
        fig.update_layout(height=300, xaxis_title="Date", yaxis_title="Hours")
        st.plotly_chart(fig, use_container_width=True)
        
        # Personnel Performance
        if activities['personnel'].notna().any():
            st.markdown("---")
            st.markdown("### 👥 Personnel Activity")
            
            # Parse personnel from comma-separated strings
            personnel_hours = {}
            for idx, row in activities.iterrows():
                if pd.notna(row['personnel']) and row['personnel']:
                    people = [p.strip() for p in row['personnel'].split(',')]
                    for person in people:
                        if person:
                            personnel_hours[person] = personnel_hours.get(person, 0) + row['hours']
            
            if personnel_hours:
                personnel_df = pd.DataFrame(list(personnel_hours.items()), columns=['Personnel', 'Hours'])
                personnel_df = personnel_df.sort_values('Hours', ascending=False)
                
                fig = px.bar(
                    personnel_df,
                    x='Personnel',
                    y='Hours',
                    color='Personnel',
                    color_discrete_sequence=['#27AE60', '#3498DB', '#9B59B6', '#E67E22', '#E74C3C', '#17A2B8', '#F8B400', '#2ECC71']
                )
                fig.update_layout(showlegend=False, height=300)
                st.plotly_chart(fig, use_container_width=True)
        
        # Contextual Insights Section
        st.markdown("---")
        st.markdown("### 💡 Key Insights & Trends")
        
        col1, col2 = st.columns(2)
        
        with col1:
            st.markdown("#### 📊 Activity Insights")
            
            # Most active day
            if not activities.empty:
                day_counts = activities['date'].value_counts()
                most_active_day = day_counts.index[0]
                st.info(f"🔥 **Busiest Day:** {pd.to_datetime(most_active_day).strftime('%A, %B %d')} ({day_counts.iloc[0]} activities)")
                
                # Average activity duration
                avg_duration = activities['hours'].mean()
                st.info(f"⏱️ **Avg Activity Duration:** {avg_duration:.2f} hours")
                
                # Most common activity type
                if activity_type_hours:
                    top_activity = max(activity_type_hours, key=activity_type_hours.get)
                    top_hours = activity_type_hours[top_activity]
                    percentage = (top_hours / sum(activity_type_hours.values())) * 100
                    st.info(f"🎯 **Top Activity:** {top_activity} ({percentage:.1f}% of hours)")
        
        with col2:
            st.markdown("#### 👥 Team Insights")
            
            if activities['personnel'].notna().any():
                # Most active personnel
                if personnel_hours:
                    top_person = max(personnel_hours, key=personnel_hours.get)
                    top_person_hours = personnel_hours[top_person]
                    st.success(f"⭐ **Top Contributor:** {top_person} ({top_person_hours:.1f} hours)")
                
                # Team utilization
                unique_personnel = len(personnel_hours) if personnel_hours else 0
                st.info(f"👥 **Active Team Members:** {unique_personnel}")
                
                # Equipment usage
                if activities['equipment'].notna().any():
                    equipment_count = 0
                    for idx, row in activities.iterrows():
                        if pd.notna(row['equipment']) and row['equipment']:
                            equipment_count += len([e.strip() for e in row['equipment'].split(',') if e.strip()])
                    st.info(f"🔧 **Equipment Uses:** {equipment_count} times")
        
        # ADVANCED VISUALIZATIONS (PHASE 2 FEATURE 3)
        st.markdown("---")
        st.markdown("### 📊 Advanced Visualizations")
        
        tab1, tab2 = st.tabs(["🔥 Activity Heatmap", "🌊 Flow Diagram"])
        
        with tab1:
            st.markdown("**Activity intensity by day and hour**")
            heatmap = create_activity_heatmap(activities)
            if heatmap:
                st.plotly_chart(heatmap, use_container_width=True)
            else:
                st.info("Not enough data with time information to generate heatmap. Add start times to activities.")
        
        with tab2:
            st.markdown("**Flow from Activity Type → Personnel → Equipment**")
            sankey = create_sankey_diagram(activities)
            if sankey:
                st.plotly_chart(sankey, use_container_width=True)
            else:
                st.info("Not enough connected data to generate flow diagram.")
    
    else:
        st.info("No activities recorded for selected period.")

#############################################
# PAGE 2.5: EXECUTIVE DASHBOARD (PHASE 2 FEATURE 5)
#############################################
elif page == "👔 Executive Dashboard":
    st.title("👔 Executive Dashboard")
    st.markdown("**High-level overview for leadership and strategic planning**")
    
    # Time period selector
    col1, col2 = st.columns([3, 1])
    with col1:
        exec_period = st.selectbox(
            "Reporting Period",
            ["This Month", "This Quarter", "This Year", "Custom"],
            key="exec_period"
        )
    
    today = datetime.now().date()
    
    if exec_period == "This Month":
        start_date = today.replace(day=1)
        end_date = today
        period_name = f"{today.strftime('%B %Y')}"
    elif exec_period == "This Quarter":
        quarter = (today.month - 1) // 3
        start_date = today.replace(month=quarter*3+1, day=1)
        end_month = quarter*3+3
        if end_month == 12:
            end_date = today.replace(month=12, day=31)
        else:
            end_date = today.replace(month=end_month+1, day=1) - timedelta(days=1)
        period_name = f"Q{quarter+1} {today.year}"
    elif exec_period == "This Year":
        start_date = today.replace(month=1, day=1)
        end_date = today
        period_name = f"{today.year}"
    else:  # Custom
        with col1:
            start_date = st.date_input("Start Date", value=today.replace(day=1), key="exec_start")
        with col2:
            end_date = st.date_input("End Date", value=today, key="exec_end")
        period_name = f"{start_date.strftime('%b %d')} - {end_date.strftime('%b %d, %Y')}"
    
    activities = get_activities(start_date, end_date)
    
    if not activities.empty:
        # EXECUTIVE SUMMARY BOX
        st.markdown("### 📊 Executive Summary")
        
        summary_col1, summary_col2, summary_col3, summary_col4 = st.columns(4)
        
        with summary_col1:
            total_hours = activities['hours'].sum()
            st.metric("Total Operating Hours", f"{total_hours:.0f}h", help="Total simulation center activity hours")
        
        with summary_col2:
            total_students = activities['students_trained'].sum()
            st.metric("Students Trained", f"{int(total_students):,}", help="Total number of students served")
        
        with summary_col3:
            utilization = (total_hours / ((end_date - start_date).days * 10)) * 100 if (end_date - start_date).days > 0 else 0
            st.metric("Utilization Rate", f"{utilization:.1f}%", help="% of available hours (assumed 10h/day)")
        
        with summary_col4:
            avg_class_size = total_students / len(activities) if len(activities) > 0 else 0
            st.metric("Avg Class Size", f"{avg_class_size:.1f}", help="Average students per activity")
        
        st.markdown("---")
        
        # AI-GENERATED EXECUTIVE NARRATIVE
        st.markdown(f"### 📖 Executive Brief - {period_name}")
        with st.spinner("Generating executive summary..."):
            exec_story = generate_data_story(activities, period_name=period_name)
            st.info(exec_story)
        
        st.markdown("---")
        
        # KEY PERFORMANCE INDICATORS
        st.markdown("### 📈 Key Performance Indicators")
        
        kpi_col1, kpi_col2, kpi_col3 = st.columns(3)
        
        with kpi_col1:
            st.markdown("#### Operational Efficiency")
            activities_per_day = len(activities) / ((end_date - start_date).days + 1)
            st.metric("Activities/Day", f"{activities_per_day:.1f}")
            
            avg_duration = activities['hours'].mean()
            st.metric("Avg Session Length", f"{avg_duration:.1f}h")
            
            completion_rate = (activities[activities['received'] > 0].shape[0] / len(activities) * 100) if len(activities) > 0 else 0
            st.metric("Turn-in Rate", f"{completion_rate:.0f}%")
        
        with kpi_col2:
            st.markdown("#### Resource Utilization")
            
            # Room utilization
            unique_rooms = set()
            for idx, row in activities.iterrows():
                if pd.notna(row['room_number']):
                    unique_rooms.update([r.strip() for r in row['room_number'].split(',')])
            st.metric("Rooms Utilized", len(unique_rooms))
            
            # Equipment diversity
            unique_equipment = set()
            for idx, row in activities.iterrows():
                if pd.notna(row['equipment']):
                    unique_equipment.update([e.strip() for e in row['equipment'].split(',')])
            st.metric("Equipment Types Used", len(unique_equipment))
            
            # Personnel engagement
            unique_personnel = set()
            for idx, row in activities.iterrows():
                if pd.notna(row['personnel']):
                    unique_personnel.update([p.strip() for p in row['personnel'].split(',')])
            st.metric("Active Staff", len(unique_personnel))
        
        with kpi_col3:
            st.markdown("#### Training Impact")
            
            unique_courses = activities['course'].dropna().nunique()
            st.metric("Courses Delivered", unique_courses)
            
            # High fidelity percentage
            hf_activities = activities[activities['activity_type'].str.contains('High Fidelity|Fidelity', case=False, na=False)]
            hf_percentage = (len(hf_activities) / len(activities) * 100) if len(activities) > 0 else 0
            st.metric("High-Fidelity %", f"{hf_percentage:.0f}%")
            
            students_per_hour = total_students / total_hours if total_hours > 0 else 0
            st.metric("Students/Hour", f"{students_per_hour:.1f}")
        
        st.markdown("---")
        
        # STRATEGIC INSIGHTS
        st.markdown("### 💡 Strategic Recommendations")
        
        insights = generate_actionable_insights(activities)
        
        if insights:
            high_impact = [i for i in insights if i['impact'] == 'high']
            medium_impact = [i for i in insights if i['impact'] == 'medium']
            
            if high_impact:
                st.markdown("#### 🔴 High Priority")
                for insight in high_impact:
                    st.error(f"{insight['icon']} **{insight['category']}**: {insight['recommendation']}")
            
            if medium_impact:
                st.markdown("#### 🟡 Medium Priority")
                for insight in medium_impact:
                    st.warning(f"{insight['icon']} **{insight['category']}**: {insight['recommendation']}")
        else:
            st.success("✅ No critical action items. Operations performing optimally.")
        
        st.markdown("---")
        
        # TREND VISUALIZATIONS FOR EXECUTIVES
        st.markdown("### 📊 Performance Trends")
        
        # Monthly trend (if enough data)
        if (end_date - start_date).days >= 30:
            monthly_data = activities.copy()
            monthly_data['month'] = pd.to_datetime(monthly_data['date']).dt.to_period('M')
            monthly_summary = monthly_data.groupby('month').agg({
                'hours': 'sum',
                'students_trained': 'sum'
            }).reset_index()
            monthly_summary['month'] = monthly_summary['month'].astype(str)
            
            fig = go.Figure()
            fig.add_trace(go.Scatter(
                x=monthly_summary['month'],
                y=monthly_summary['hours'],
                name='Hours',
                line=dict(color='#F8B400', width=3),
                fill='tozeroy'
            ))
            fig.update_layout(
                title="Monthly Operating Hours Trend",
                xaxis_title="Month",
                yaxis_title="Hours",
                height=300
            )
            st.plotly_chart(fig, use_container_width=True)
        
        # Activity Type Distribution (Executive View)
        col1, col2 = st.columns(2)
        
        with col1:
            activity_type_hours = {}
            for idx, row in activities.iterrows():
                if pd.notna(row['activity_type']):
                    types = [t.strip() for t in row['activity_type'].split(',')]
                    for act_type in types:
                        activity_type_hours[act_type] = activity_type_hours.get(act_type, 0) + row['hours']
            
            if activity_type_hours:
                fig = px.pie(
                    values=list(activity_type_hours.values()),
                    names=list(activity_type_hours.keys()),
                    title="Activity Type Distribution",
                    color_discrete_sequence=['#F8B400', '#000000', '#4A90E2', '#27AE60', '#E67E22']
                )
                fig.update_traces(textposition='inside', textinfo='percent+label')
                fig.update_layout(height=350, showlegend=False)
                st.plotly_chart(fig, use_container_width=True)
        
        with col2:
            # Top Courses Bar Chart
            top_courses = activities[activities['course'].notna()].groupby('course')['students_trained'].sum().sort_values(ascending=False).head(5)
            
            if not top_courses.empty:
                fig = px.bar(
                    x=top_courses.values,
                    y=top_courses.index,
                    orientation='h',
                    title="Top 5 Courses by Students",
                    labels={'x': 'Students', 'y': 'Course'},
                    color_discrete_sequence=['#F8B400']
                )
                fig.update_layout(height=350, showlegend=False)
                st.plotly_chart(fig, use_container_width=True)
        
        st.markdown("---")
        
        # EXPORTABLE SUMMARY FOR REPORTS
        st.markdown("### 📄 Export Options")
        
        export_data = {
            'Period': [period_name],
            'Total Hours': [f"{total_hours:.0f}"],
            'Students Trained': [int(total_students)],
            'Activities': [len(activities)],
            'Courses': [unique_courses],
            'Utilization Rate': [f"{utilization:.1f}%"],
            'Avg Class Size': [f"{avg_class_size:.1f}"]
        }
        
        export_df = pd.DataFrame(export_data)
        
        # Check if PowerPoint is available
        try:
            import pptx
            pptx_available = True
        except ImportError:
            pptx_available = False
        
        if pptx_available:
            col1, col2, col3 = st.columns(3)
        else:
            col1, col3 = st.columns(2)
        
        with col1:
            csv = export_df.to_csv(index=False)
            st.download_button(
                label="📥 Download CSV",
                data=csv,
                file_name=f"Executive_Summary_{period_name.replace(' ', '_')}.csv",
                mime="text/csv",
                use_container_width=True
            )
        
        if pptx_available:
            with col2:
                # PowerPoint Export
                if st.button("📊 Generate PowerPoint", type="primary", use_container_width=True):
                    with st.spinner("Creating PowerPoint presentation..."):
                        pptx_data, error = create_powerpoint_report(activities, period_name, start_date, end_date)
                        
                        if error:
                            st.error(error)
                        elif pptx_data:
                            st.success("✅ PowerPoint created successfully!")
                            st.download_button(
                                label="📥 Download PowerPoint",
                                data=pptx_data,
                                file_name=f"VCU_SimCenter_Briefing_{period_name.replace(' ', '_')}.pptx",
                                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                                type="primary",
                                use_container_width=True
                            )
        
        with col3:
            if pptx_available:
                st.info("💡 **PowerPoint includes:**\n- Executive Summary\n- Key Metrics\n- AI Insights\n- Activity Breakdown\n- Recommendations\n- Top Courses")
            else:
                st.warning("📊 **PowerPoint Export Disabled**\n\nTo enable:\n```\npip install python-pptx --break-system-packages\n```\nThen restart tracker")
    
    else:
        st.info("No activities recorded for selected period.")

#############################################
# PAGE 3: COURSE ANALYTICS (COURSE-FOCUSED)
#############################################
elif page == "📚 Course Analytics":
    st.title("📚 Course Analytics & Scheduling")
    
    # Date range selector
    col1, col2, col3, col4 = st.columns(4)
    with col1:
        date_range = st.selectbox(
            "Analysis Period",
            ["All Time", "This Month", "Last Month", "Last 3 Months", "Custom"],
            key="course_range"
        )
    
    with col2:
        min_hours_filter = st.number_input("Min Hours", min_value=0.0, value=0.0, step=0.5, key="course_min_hours")
    
    with col3:
        min_sessions_filter = st.number_input("Min Sessions", min_value=0, value=0, step=1, key="course_min_sessions")
    
    with col4:
        sort_by = st.selectbox("Sort By", ["Total Hours", "Session Count", "Students", "Most Recent"], key="course_sort")
    
    today = datetime.now().date()
    
    if date_range == "All Time":
        start_date = None
        end_date = None
    elif date_range == "This Month":
        start_date = today.replace(day=1)
        end_date = today
    elif date_range == "Last Month":
        first_this_month = today.replace(day=1)
        end_date = first_this_month - timedelta(days=1)
        start_date = end_date.replace(day=1)
    elif date_range == "Last 3 Months":
        start_date = today - timedelta(days=90)
        end_date = today
    else:  # Custom
        with col2:
            start_date = st.date_input("Start", value=today - timedelta(days=90), key="course_start")
        with col3:
            end_date = st.date_input("End", value=today, key="course_end")
    
    # Get course analytics data
    conn = sqlite3.connect(DB_PATH)
    
    if start_date and end_date:
        date_filter = f"AND date BETWEEN '{start_date}' AND '{end_date}'"
    else:
        date_filter = ""
    
    course_query = f"""
        SELECT 
            course,
            COUNT(*) as session_count,
            SUM(hours) as total_hours,
            AVG(hours) as avg_hours_per_session,
            MIN(date) as first_session,
            MAX(date) as last_session
        FROM activities
        WHERE course IS NOT NULL AND course != '' {date_filter}
        GROUP BY course
        ORDER BY total_hours DESC
    """
    
    course_analytics = pd.read_sql_query(course_query, conn)
    
    # Apply filters
    if not course_analytics.empty:
        if min_hours_filter > 0:
            course_analytics = course_analytics[course_analytics['total_hours'] >= min_hours_filter]
        
        if min_sessions_filter > 0:
            course_analytics = course_analytics[course_analytics['session_count'] >= min_sessions_filter]
        
        # Apply sorting
        if sort_by == "Total Hours":
            course_analytics = course_analytics.sort_values('total_hours', ascending=False)
        elif sort_by == "Session Count":
            course_analytics = course_analytics.sort_values('session_count', ascending=False)
        elif sort_by == "Most Recent":
            course_analytics = course_analytics.sort_values('last_session', ascending=False)
    
    # Keep connection open for other tabs
    
    if not course_analytics.empty:
        # Top-level metrics
        col1, col2, col3, col4 = st.columns(4)
        
        with col1:
            total_courses = len(course_analytics)
            st.metric("📚 Active Courses", total_courses)
        
        with col2:
            total_sessions = course_analytics['session_count'].sum()
            st.metric("📅 Total Sessions", int(total_sessions))
        
        with col3:
            total_course_hours = course_analytics['total_hours'].sum()
            st.metric("⏰ Total Course Hours", f"{total_course_hours:.1f}")
        
        with col4:
            avg_session_length = course_analytics['avg_hours_per_session'].mean()
            st.metric("⌛ Avg Session Length", f"{avg_session_length:.2f}h")
        
        # Show filter info if filters are active
        if min_hours_filter > 0 or min_sessions_filter > 0:
            filter_info = []
            if min_hours_filter > 0:
                filter_info.append(f"≥{min_hours_filter}h")
            if min_sessions_filter > 0:
                filter_info.append(f"≥{min_sessions_filter} sessions")
            st.info(f"🔍 Filters active: {', '.join(filter_info)}")
        
        st.markdown("---")
        
        # Visualization section
        tab1, tab2, tab3, tab4 = st.tabs(["📊 Overview", "🏆 Top Courses", "📈 Trends", "📋 Course Details"])
        
        with tab1:
            col1, col2 = st.columns(2)
            
            with col1:
                st.markdown("### Hours by Course")
                fig = px.bar(
                    course_analytics.head(10),
                    x='total_hours',
                    y='course',
                    orientation='h',
                    labels={'total_hours': 'Total Hours', 'course': 'Course'},
                    color='course',
                    color_discrete_sequence=['#9B59B6', '#3498DB', '#27AE60', '#E67E22', '#E74C3C', '#17A2B8', '#F8B400', '#1ABC9C', '#E91E63', '#FF9800']
                )
                fig.update_layout(showlegend=False, height=400, yaxis={'categoryorder':'total ascending'})
                st.plotly_chart(fig, use_container_width=True)
            
            with col2:
                st.markdown("### Session Count by Course")
                fig = px.bar(
                    course_analytics.head(10),
                    x='session_count',
                    y='course',
                    orientation='h',
                    labels={'session_count': 'Sessions', 'course': 'Course'},
                    color='course',
                    color_discrete_sequence=['#17A2B8', '#27AE60', '#9B59B6', '#E67E22', '#E74C3C', '#3498DB', '#F8B400', '#1ABC9C', '#E91E63', '#FF9800']
                )
                fig.update_layout(showlegend=False, height=400, yaxis={'categoryorder':'total ascending'})
                st.plotly_chart(fig, use_container_width=True)
        
        with tab2:
            st.markdown("### 🏆 Top Performing Courses")
            
            col1, col2 = st.columns(2)
            
            with col1:
                st.markdown("#### Most Hours Delivered")
                top_hours = course_analytics.nlargest(5, 'total_hours')[['course', 'total_hours', 'session_count']]
                for idx, row in top_hours.iterrows():
                    st.markdown(f"""
                    <div style='background: linear-gradient(135deg, #F8B400 0%, #000000 100%); 
                                padding: 15px; border-radius: 10px; margin: 10px 0; color: white;'>
                        <h4 style='margin:0; color: white;'>{row['course']}</h4>
                        <p style='margin:5px 0;'>⏰ {row['total_hours']:.1f} hours • 📅 {int(row['session_count'])} sessions</p>
                    </div>
                    """, unsafe_allow_html=True)
            
            with col2:
                st.markdown("#### Most Frequently Run")
                top_frequent = course_analytics.nlargest(5, 'session_count')[['course', 'session_count', 'avg_hours_per_session']]
                for idx, row in top_frequent.iterrows():
                    st.markdown(f"""
                    <div style='background: linear-gradient(135deg, #000000 0%, #F8B400 100%); 
                                padding: 15px; border-radius: 10px; margin: 10px 0; color: white;'>
                        <h4 style='margin:0; color: white;'>{row['course']}</h4>
                        <p style='margin:5px 0;'>📅 {int(row['session_count'])} sessions • ⌛ {row['avg_hours_per_session']:.2f}h avg</p>
                    </div>
                    """, unsafe_allow_html=True)
            
            st.markdown("---")
            
            # Course efficiency metrics
            st.markdown("### 📈 Course Metrics")
            col1, col2, col3 = st.columns(3)
            
            with col1:
                # Longest average sessions
                longest_sessions = course_analytics.nlargest(3, 'avg_hours_per_session')[['course', 'avg_hours_per_session']]
                st.markdown("**Longest Sessions**")
                for idx, row in longest_sessions.iterrows():
                    st.write(f"• {row['course']}: {row['avg_hours_per_session']:.2f}h")
            
            with col2:
                # Most consistent courses (highest session count)
                most_consistent = course_analytics.nlargest(3, 'session_count')[['course', 'session_count']]
                st.markdown("**Most Consistent**")
                for idx, row in most_consistent.iterrows():
                    st.write(f"• {row['course']}: {int(row['session_count'])} runs")
            
            with col3:
                # FIXED: Convert to datetime first
                course_analytics['last_session_dt'] = pd.to_datetime(course_analytics['last_session'])
                recent = course_analytics.nlargest(3, 'last_session_dt')[['course', 'last_session']]
                st.markdown("**Most Recent Activity**")
                for idx, row in recent.iterrows():
                    st.write(f"• {row['course']}: {row['last_session']}")
        
        with tab3:
            st.markdown("### 📈 Course Activity Trends")
            
            # Get time-series data for courses (reuse existing connection)
            trend_query = f"""
                SELECT 
                    date,
                    course,
                    SUM(hours) as hours
                FROM activities
                WHERE course IS NOT NULL AND course != '' {date_filter}
                GROUP BY date, course
                ORDER BY date
            """
            
            trend_data = pd.read_sql_query(trend_query, conn)
            
            if not trend_data.empty:
                # Select top 5 courses to visualize
                top_courses = course_analytics.head(5)['course'].tolist()
                trend_filtered = trend_data[trend_data['course'].isin(top_courses)]
                
                st.markdown("#### Course Hours Over Time (Top 5 Courses)")
                fig = px.line(
                    trend_filtered,
                    x='date',
                    y='hours',
                    color='course',
                    markers=True,
                    labels={'hours': 'Hours', 'date': 'Date', 'course': 'Course'}
                )
                fig.update_layout(height=400)
                st.plotly_chart(fig, use_container_width=True)
                
                # Monthly aggregation
                st.markdown("---")
                st.markdown("#### Monthly Course Activity Heatmap")
                
                monthly_query = f"""
                    SELECT 
                        strftime('%Y-%m', date) as month,
                        course,
                        SUM(hours) as total_hours
                    FROM activities
                    WHERE course IS NOT NULL AND course != '' {date_filter}
                    GROUP BY month, course
                    ORDER BY month DESC
                """
                
                monthly_data = pd.read_sql_query(monthly_query, conn)
                
                if not monthly_data.empty:
                    # Pivot for heatmap-style visualization
                    pivot_hours = monthly_data.pivot_table(
                        index='course',
                        columns='month',
                        values='total_hours',
                        fill_value=0
                    )
                    
                    fig = px.imshow(
                        pivot_hours,
                        labels=dict(x="Month", y="Course", color="Hours"),
                        color_continuous_scale=['white', '#F8B400', '#000000'],
                        aspect="auto"
                    )
                    fig.update_layout(height=600)
                    st.plotly_chart(fig, use_container_width=True)
                
                # Day of week analysis
                st.markdown("---")
                st.markdown("#### Course Activity by Day of Week")
                
                trend_data['DayOfWeek'] = pd.to_datetime(trend_data['date']).dt.day_name()
                day_order = ['Monday', 'Tuesday', 'Wednesday', 'Thursday', 'Friday', 'Saturday', 'Sunday']
                
                dow_hours = trend_data.groupby('DayOfWeek')['hours'].sum().reset_index()
                dow_hours['DayOfWeek'] = pd.Categorical(dow_hours['DayOfWeek'], categories=day_order, ordered=True)
                dow_hours = dow_hours.sort_values('DayOfWeek')
                
                fig = px.bar(
                    dow_hours,
                    x='DayOfWeek',
                    y='hours',
                    labels={'DayOfWeek': 'Day of Week', 'hours': 'Total Hours'},
                    color='DayOfWeek',
                    color_discrete_sequence=['#E67E22', '#3498DB', '#27AE60', '#9B59B6', '#E74C3C', '#17A2B8', '#F8B400']
                )
                fig.update_layout(showlegend=False, height=300)
                st.plotly_chart(fig, use_container_width=True)
        
        with tab4:
            st.markdown("### 📋 Course Schedule & Details")
            
            # Get course details with scheduling info
            courses_df = get_courses()
            
            # Just show analytics (courses table is simplified now - no room/time data there)
            display_df = course_analytics.copy()
            display_df['total_hours'] = display_df['total_hours'].round(2)
            display_df['avg_hours_per_session'] = display_df['avg_hours_per_session'].round(2)
            
            display_df = display_df.rename(columns={
                'course': 'Course',
                'session_count': 'Sessions',
                'total_hours': 'Total Hours',
                'avg_hours_per_session': 'Avg Hours/Session',
                'first_session': 'First Run',
                'last_session': 'Last Run'
            })
            
            st.dataframe(display_df, use_container_width=True, hide_index=True)
            
            # Download button
            csv = display_df.to_csv(index=False)
            st.download_button(
                label="📥 Download Course Analytics CSV",
                data=csv,
                file_name=f"course_analytics_{datetime.now().strftime('%Y%m%d')}.csv",
                mime="text/csv"
            )
            
            # Close connection after all queries done
            if 'conn' in locals():
                conn.close()
    
    else:
        st.info("No course data available for the selected period.")
    
    # Final connection cleanup
    if 'conn' in locals() and not conn:
        try:
            conn.close()
        except:
            pass

#############################################
# PAGE 4: EQUIPMENT ANALYTICS (ERROR-FREE)
#############################################
elif page == "🔧 Equipment Analytics":
    st.title("🔧 Equipment Analytics & Utilization")
    
    # Date range selector
    col1, col2, col3 = st.columns(3)
    with col1:
        date_range = st.selectbox(
            "Analysis Period",
            ["All Time", "This Month", "Last Month", "Last 3 Months", "Custom"],
            key="equipment_range"
        )
    
    today = datetime.now().date()
    
    if date_range == "All Time":
        start_date = None
        end_date = None
    elif date_range == "This Month":
        start_date = today.replace(day=1)
        end_date = today
    elif date_range == "Last Month":
        first_this_month = today.replace(day=1)
        end_date = first_this_month - timedelta(days=1)
        start_date = end_date.replace(day=1)
    elif date_range == "Last 3 Months":
        start_date = today - timedelta(days=90)
        end_date = today
    else:  # Custom
        with col2:
            start_date = st.date_input("Start", value=today - timedelta(days=90), key="equipment_start")
        with col3:
            end_date = st.date_input("End", value=today, key="equipment_end")
    
    # Get all activities with equipment
    if start_date and end_date:
        activities = get_activities(start_date, end_date)
    else:
        activities = get_activities()
    
    # Parse equipment usage (equipment is stored as comma-separated strings)
    equipment_usage = {}
    equipment_hours = {}
    equipment_dates = {}
    
    for idx, row in activities.iterrows():
        if pd.notna(row['equipment']) and row['equipment']:
            equipment_items = [e.strip() for e in row['equipment'].split(',')]
            for equipment in equipment_items:
                if equipment:
                    equipment_usage[equipment] = equipment_usage.get(equipment, 0) + 1
                    equipment_hours[equipment] = equipment_hours.get(equipment, 0) + row['hours']
                    
                    if equipment not in equipment_dates:
                        equipment_dates[equipment] = []
                    equipment_dates[equipment].append(row['date'])
    
    if equipment_usage:
        # Create comprehensive analytics dataframe
        equipment_analytics = pd.DataFrame({
            'Equipment': list(equipment_usage.keys()),
            'Usage Count': list(equipment_usage.values()),
            'Total Hours': [equipment_hours[e] for e in equipment_usage.keys()],
            'Avg Hours/Use': [equipment_hours[e]/equipment_usage[e] for e in equipment_usage.keys()],
            'Last Used': [max(equipment_dates[e]) for e in equipment_usage.keys()]
        })
        
        equipment_analytics = equipment_analytics.sort_values('Usage Count', ascending=False)
        
        # Get equipment status information
        equipment_status_df = get_equipment()
        
        # Top-level metrics
        col1, col2, col3, col4 = st.columns(4)
        
        with col1:
            total_equipment_tracked = len(equipment_analytics)
            st.metric("🔧 Equipment Items Tracked", total_equipment_tracked)
        
        with col2:
            total_uses = equipment_analytics['Usage Count'].sum()
            st.metric("📊 Total Uses", int(total_uses))
        
        with col3:
            total_eq_hours = equipment_analytics['Total Hours'].sum()
            st.metric("⏰ Total Equipment Hours", f"{total_eq_hours:.1f}")
        
        with col4:
            avg_uses_per_item = equipment_analytics['Usage Count'].mean()
            st.metric("📈 Avg Uses/Item", f"{avg_uses_per_item:.1f}")
        
        st.markdown("---")
        
        # Tabs for different views
        tab1, tab2, tab3, tab4 = st.tabs(["📊 Usage Overview", "🏆 Top Equipment", "📅 Utilization Trends", "📋 Detailed Table"])
        
        with tab1:
            col1, col2 = st.columns(2)
            
            with col1:
                st.markdown("### Most Used Equipment (by count)")
                fig = px.bar(
                    equipment_analytics.head(10),
                    x='Usage Count',
                    y='Equipment',
                    orientation='h',
                    labels={'Usage Count': 'Times Used', 'Equipment': 'Equipment'},
                    color='Equipment',
                    color_discrete_sequence=['#3498DB', '#27AE60', '#9B59B6', '#E67E22', '#E74C3C', '#17A2B8', '#F8B400', '#1ABC9C', '#E91E63', '#FF9800']
                )
                fig.update_layout(showlegend=False, height=400, yaxis={'categoryorder':'total ascending'})
                st.plotly_chart(fig, use_container_width=True)
            
            with col2:
                st.markdown("### Equipment Hours Distribution")
                fig = px.bar(
                    equipment_analytics.head(10),
                    x='Total Hours',
                    y='Equipment',
                    orientation='h',
                    labels={'Total Hours': 'Hours', 'Equipment': 'Equipment'},
                    color='Equipment',
                    color_discrete_sequence=['#1ABC9C', '#E74C3C', '#3498DB', '#F8B400', '#9B59B6', '#27AE60', '#E67E22', '#17A2B8', '#E91E63', '#FF9800']
                )
                fig.update_layout(showlegend=False, height=400, yaxis={'categoryorder':'total ascending'})
                st.plotly_chart(fig, use_container_width=True)
            
            st.markdown("---")
            
            # Equipment status overview
            st.markdown("### 🔍 Equipment Status Overview")
            
            if not equipment_status_df.empty:
                col1, col2, col3, col4 = st.columns(4)
                
                status_counts = equipment_status_df['status'].value_counts()
                
                with col1:
                    operational_count = status_counts.get('Operational', 0)
                    st.metric("✅ Operational", operational_count)
                with col2:
                    maintenance_count = status_counts.get('Maintenance', 0)
                    st.metric("🔧 Maintenance", maintenance_count)
                with col3:
                    down_count = status_counts.get('Down', 0)
                    st.metric("❌ Down", down_count)
                with col4:
                    total = len(equipment_status_df)
                    uptime = (operational_count / total * 100) if total > 0 else 0
                    st.metric("📈 Uptime Rate", f"{uptime:.1f}%")
            else:
                st.info("Add equipment in Settings to track status.")
        
        with tab2:
            st.markdown("### 🏆 Top Performing Equipment")
            
            col1, col2 = st.columns(2)
            
            with col1:
                st.markdown("#### Most Frequently Used")
                top_usage = equipment_analytics.nlargest(5, 'Usage Count')[['Equipment', 'Usage Count', 'Total Hours']]
                for idx, row in top_usage.iterrows():
                    st.markdown(f"""
                    <div style='background: linear-gradient(135deg, #F8B400 0%, #000000 100%); 
                                padding: 15px; border-radius: 10px; margin: 10px 0; color: white;'>
                        <h4 style='margin:0; color: white;'>{row['Equipment']}</h4>
                        <p style='margin:5px 0;'>🔄 {int(row['Usage Count'])} uses • ⏰ {row['Total Hours']:.1f} hours</p>
                    </div>
                    """, unsafe_allow_html=True)
            
            with col2:
                st.markdown("#### Longest Operating Hours")
                top_hours = equipment_analytics.nlargest(5, 'Total Hours')[['Equipment', 'Total Hours', 'Usage Count']]
                for idx, row in top_hours.iterrows():
                    st.markdown(f"""
                    <div style='background: linear-gradient(135deg, #000000 0%, #F8B400 100%); 
                                padding: 15px; border-radius: 10px; margin: 10px 0; color: white;'>
                        <h4 style='margin:0; color: white;'>{row['Equipment']}</h4>
                        <p style='margin:5px 0;'>⏰ {row['Total Hours']:.1f} hours • 🔄 {int(row['Usage Count'])} uses</p>
                    </div>
                    """, unsafe_allow_html=True)
            
            st.markdown("---")
            
            # Utilization insights
            st.markdown("### 📊 Utilization Insights")
            
            col1, col2, col3 = st.columns(3)
            
            with col1:
                st.markdown("**Average Session Length**")
                longest_sessions = equipment_analytics.nlargest(3, 'Avg Hours/Use')[['Equipment', 'Avg Hours/Use']]
                for idx, row in longest_sessions.iterrows():
                    st.write(f"• {row['Equipment']}: {row['Avg Hours/Use']:.2f}h")
            
            with col2:
                # Calculate utilization rate (uses per day in period)
                if start_date and end_date:
                    days_in_period = (end_date - start_date).days + 1
                else:
                    days_in_period = 365  # Assume 1 year for all time
                
                equipment_analytics['Uses Per Day'] = equipment_analytics['Usage Count'] / days_in_period
                high_utilization = equipment_analytics.nlargest(3, 'Uses Per Day')[['Equipment', 'Uses Per Day']]
                st.markdown("**Highest Utilization Rate**")
                for idx, row in high_utilization.iterrows():
                    st.write(f"• {row['Equipment']}: {row['Uses Per Day']:.3f}/day")
            
            with col3:
                # FIXED: Convert to datetime first
                equipment_analytics['Last Used Date'] = pd.to_datetime(equipment_analytics['Last Used'])
                recent = equipment_analytics.nlargest(3, 'Last Used Date')[['Equipment', 'Last Used']]
                st.markdown("**Most Recently Used**")
                for idx, row in recent.iterrows():
                    st.write(f"• {row['Equipment']}: {row['Last Used']}")
            
            st.markdown("---")
            
            # Underutilized equipment
            if len(equipment_analytics) >= 5:
                st.markdown("### ⚠️ Underutilized Equipment")
                underutilized = equipment_analytics.nsmallest(5, 'Usage Count')[['Equipment', 'Usage Count', 'Last Used']]
                
                st.dataframe(
                    underutilized,
                    use_container_width=True,
                    hide_index=True
                )
                
                st.info("💡 Consider reviewing training needs or equipment placement for items with low utilization.")
        
        with tab3:
            st.markdown("### 📅 Equipment Utilization Trends")
            
            # Timeline visualization for top equipment
            top_equipment = equipment_analytics.head(5)['Equipment'].tolist()
            
            timeline_data = []
            for equipment in top_equipment:
                for date in equipment_dates[equipment]:
                    timeline_data.append({
                        'Equipment': equipment,
                        'Date': date,
                        'Count': 1
                    })
            
            if timeline_data:
                timeline_df = pd.DataFrame(timeline_data)
                
                # Aggregate by date
                timeline_agg = timeline_df.groupby(['Date', 'Equipment']).sum().reset_index()
                
                st.markdown("#### Usage Over Time (Top 5 Equipment)")
                fig = px.line(
                    timeline_agg,
                    x='Date',
                    y='Count',
                    color='Equipment',
                    markers=True,
                    labels={'Count': 'Times Used', 'Date': 'Date', 'Equipment': 'Equipment'}
                )
                fig.update_layout(height=400)
                st.plotly_chart(fig, use_container_width=True)
                
                # Monthly heatmap
                st.markdown("---")
                st.markdown("#### Monthly Usage Heatmap")
                
                timeline_df['Month'] = pd.to_datetime(timeline_df['Date']).dt.strftime('%Y-%m')
                monthly = timeline_df.groupby(['Month', 'Equipment']).size().reset_index(name='Uses')
                
                pivot_monthly = monthly.pivot_table(
                    index='Equipment',
                    columns='Month',
                    values='Uses',
                    fill_value=0
                )
                
                fig = px.imshow(
                    pivot_monthly,
                    labels=dict(x="Month", y="Equipment", color="Uses"),
                    color_continuous_scale=['white', '#F8B400', '#000000'],
                    aspect="auto"
                )
                fig.update_layout(height=500)
                st.plotly_chart(fig, use_container_width=True)
                
                # Day of week analysis
                st.markdown("---")
                st.markdown("#### Usage by Day of Week")
                
                timeline_df['DayOfWeek'] = pd.to_datetime(timeline_df['Date']).dt.day_name()
                day_order = ['Monday', 'Tuesday', 'Wednesday', 'Thursday', 'Friday', 'Saturday', 'Sunday']
                
                dow_usage = timeline_df.groupby('DayOfWeek').size().reset_index(name='Uses')
                dow_usage['DayOfWeek'] = pd.Categorical(dow_usage['DayOfWeek'], categories=day_order, ordered=True)
                dow_usage = dow_usage.sort_values('DayOfWeek')
                
                fig = px.bar(
                    dow_usage,
                    x='DayOfWeek',
                    y='Uses',
                    labels={'DayOfWeek': 'Day of Week', 'Uses': 'Total Uses'},
                    color='DayOfWeek',
                    color_discrete_sequence=['#E74C3C', '#3498DB', '#27AE60', '#F8B400', '#9B59B6', '#17A2B8', '#E67E22']
                )
                fig.update_layout(showlegend=False, height=300)
                st.plotly_chart(fig, use_container_width=True)
        
        with tab4:
            st.markdown("### 📋 Complete Equipment Analytics Table")
            
            # Display dataframe
            display_df = equipment_analytics.copy()
            display_df['Total Hours'] = display_df['Total Hours'].round(2)
            display_df['Avg Hours/Use'] = display_df['Avg Hours/Use'].round(2)
            
            st.dataframe(
                display_df,
                use_container_width=True,
                hide_index=True
            )
            
            # Download button
            csv = display_df.to_csv(index=False)
            st.download_button(
                label="📥 Download Equipment Analytics CSV",
                data=csv,
                file_name=f"equipment_analytics_{datetime.now().strftime('%Y%m%d')}.csv",
                mime="text/csv"
            )
            
            st.markdown("---")
            
            # Equipment recommendations
            st.markdown("### 💡 Recommendations")
            
            col1, col2 = st.columns(2)
            
            with col1:
                st.markdown("**🔧 Maintenance Priority**")
                # Identify high-use equipment that may need maintenance attention
                high_use = equipment_analytics.nlargest(3, 'Usage Count')
                for idx, row in high_use.iterrows():
                    days_since_last = (today - pd.to_datetime(row['Last Used']).date()).days
                    st.write(f"• **{row['Equipment']}**: {int(row['Usage Count'])} uses, last used {days_since_last} days ago")
            
            with col2:
                st.markdown("**📊 Low Utilization Items**")
                # Identify underutilized equipment
                if len(equipment_analytics) >= 3:
                    low_use = equipment_analytics.nsmallest(3, 'Usage Count')
                    for idx, row in low_use.iterrows():
                        st.write(f"• **{row['Equipment']}**: Only {int(row['Usage Count'])} uses")
                    st.info("Consider promoting training sessions for these items")
    
    else:
        st.info("No equipment usage data available for the selected period.")

# [CONTINUING WITH REMAINING PAGES...]
#############################################
# PAGE 5: HISTORY WITH EDITING
#############################################
elif page == "📅 History":
    st.title("📅 Activity History")
    
    col1, col2 = st.columns(2)
    with col1:
        start_date = st.date_input("Start Date", value=datetime.now() - timedelta(days=30))
    with col2:
        end_date = st.date_input("End Date", value=datetime.now())
    
    activities = get_activities(start_date, end_date)
    
    if not activities.empty:
        st.markdown(f"### Showing {len(activities)} activities")
        
        # Filter options
        col1, col2, col3 = st.columns(3)
        
        with col1:
            activity_filter = st.multiselect(
                "Filter by Activity Type",
                activities['activity_type'].unique()
            )
        
        with col2:
            course_filter = st.multiselect(
                "Filter by Course",
                activities['course'].dropna().unique()
            )
        
        with col3:
            search_term = st.text_input("Search in notes")
        
        # Apply filters
        filtered = activities.copy()
        
        if activity_filter:
            filtered = filtered[filtered['activity_type'].isin(activity_filter)]
        
        if course_filter:
            filtered = filtered[filtered['course'].isin(course_filter)]
        
        if search_term:
            filtered = filtered[filtered['notes'].str.contains(search_term, case=False, na=False)]
        
        st.markdown("---")
        
        for idx, row in filtered.iterrows():
            with st.expander(f"📅 {row['date']} - {row['activity_type']} ({row['hours']}h)"):
                # View mode
                if f"edit_{row['id']}" not in st.session_state:
                    col1, col2 = st.columns(2)
                    
                    with col1:
                        st.write(f"**Activity Type:** {row['activity_type']}")
                        st.write(f"**Hours:** {row['hours']}")
                        st.write(f"**Students Trained:** {row['students_trained']}")
                        st.write(f"**Personnel:** {row['personnel'] or 'N/A'}")
                    
                    with col2:
                        st.write(f"**Course:** {row['course'] or 'N/A'}")
                        st.write(f"**Equipment:** {row['equipment'] or 'N/A'}")
                        st.write(f"**Turn In:** {row['turn_in']}")
                        st.write(f"**Received:** {row['received']}")
                    
                    # Display room and time info
                    if pd.notna(row.get('room_number')) and row['room_number']:
                        st.write(f"**Room(s):** {row['room_number']}")
                    if pd.notna(row.get('time_start')) and row['time_start']:
                        time_display = f"{row['time_start']}"
                        if pd.notna(row.get('time_end')) and row['time_end']:
                            time_display += f" - {row['time_end']}"
                        st.write(f"**Time:** {time_display}")
                    
                    if row['notes']:
                        st.write(f"**Notes:** {row['notes']}")
                    
                    col1, col2 = st.columns(2)
                    with col1:
                        if st.button("✏️ Edit", key=f"edit_btn_{row['id']}"):
                            st.session_state[f"edit_{row['id']}"] = True
                            st.rerun()
                    with col2:
                        if st.button("🗑️ Delete", key=f"del_{row['id']}"):
                            delete_activity(row['id'])
                            st.success("Activity deleted")
                            st.rerun()
                
                # Edit mode
                else:
                    st.markdown("**✏️ Editing Activity**")
                    
                    edit_date = st.date_input("Date", value=pd.to_datetime(row['date']).date(), key=f"edate_{row['id']}")
                    
                    activity_types = get_activity_types()
                    available_types = activity_types['name'].tolist()
                    
                    # Parse existing activity types and validate they exist
                    existing_types = []
                    if pd.notna(row['activity_type']) and row['activity_type']:
                        parsed_types = [t.strip() for t in row['activity_type'].split(',')]
                        # Only include types that are in the available list
                        existing_types = [t for t in parsed_types if t in available_types]
                    
                    edit_activity_types = st.multiselect(
                        "Activity Type(s)",
                        available_types,
                        default=existing_types,
                        key=f"etype_{row['id']}"
                    )
                    
                    col1, col2 = st.columns(2)
                    with col1:
                        edit_hours = st.number_input("Hours", min_value=0.0, max_value=24.0, step=0.25, value=float(row['hours']), key=f"ehrs_{row['id']}")
                        edit_students = st.number_input("Students", min_value=0, step=1, value=int(row['students_trained']), key=f"estu_{row['id']}")
                    with col2:
                        edit_turn_in = st.number_input("Turn In", min_value=0, step=1, value=int(row['turn_in']), key=f"etin_{row['id']}")
                        edit_received = st.number_input("Received", min_value=0, step=1, value=int(row['received']), key=f"erec_{row['id']}")
                    
                    courses_df = get_courses()
                    course_list = [""] + courses_df['name'].tolist()
                    edit_course = st.selectbox(
                        "Course",
                        course_list,
                        index=course_list.index(row['course']) if row['course'] in course_list else 0,
                        key=f"ecrs_{row['id']}"
                    )
                    
                    # Room and Time fields
                    col1, col2, col3 = st.columns(3)
                    with col1:
                        # Parse existing rooms (may be comma-separated)
                        existing_rooms = []
                        if pd.notna(row.get('room_number')) and row['room_number']:
                            existing_rooms = [r.strip() for r in str(row['room_number']).split(',')]
                        edit_rooms = st.multiselect("Room(s)", ROOM_NUMBERS, 
                                                     default=existing_rooms,
                                                     key=f"eroom_{row['id']}")
                    with col2:
                        if pd.notna(row.get('time_start')) and row['time_start']:
                            default_start = datetime.strptime(row['time_start'], "%H:%M").time()
                        else:
                            default_start = None
                        edit_time_start = st.time_input("Start Time", value=default_start, key=f"ets_{row['id']}")
                    with col3:
                        if pd.notna(row.get('time_end')) and row['time_end']:
                            default_end = datetime.strptime(row['time_end'], "%H:%M").time()
                        else:
                            default_end = None
                        edit_time_end = st.time_input("End Time", value=default_end, key=f"ete_{row['id']}")
                    
                    edit_personnel = st.text_input("Personnel (comma-separated)", value=row['personnel'] or "", key=f"eper_{row['id']}")
                    edit_equipment = st.text_input("Equipment (comma-separated)", value=row['equipment'] or "", key=f"eeq_{row['id']}")
                    edit_notes = st.text_area("Notes", value=row['notes'] or "", key=f"enot_{row['id']}")
                    
                    col1, col2 = st.columns(2)
                    with col1:
                        if st.button("💾 Save Changes", key=f"save_{row['id']}", type="primary"):
                            edit_time_start_str = edit_time_start.strftime("%H:%M") if edit_time_start else ""
                            edit_time_end_str = edit_time_end.strftime("%H:%M") if edit_time_end else ""
                            edit_rooms_str = ", ".join(edit_rooms) if edit_rooms else ""
                            edit_activity_types_str = ", ".join(edit_activity_types) if edit_activity_types else ""
                            
                            update_activity(
                                row['id'], edit_date, edit_activity_types_str, edit_hours, edit_students,
                                edit_personnel, edit_equipment, edit_course, edit_rooms_str,
                                edit_time_start_str, edit_time_end_str, edit_turn_in, edit_received, edit_notes
                            )
                            del st.session_state[f"edit_{row['id']}"]
                            st.success("Activity updated!")
                            st.rerun()
                    with col2:
                        if st.button("❌ Cancel", key=f"cancel_{row['id']}"):
                            del st.session_state[f"edit_{row['id']}"]
                            st.rerun()
    else:
        st.info("No activities found for selected date range.")

#############################################
# PAGE 5.5: CANCELLATIONS TRACKER (NEW!)
#############################################
elif page == "❌ Cancellations":
    st.title("❌ Course Cancellations Tracker")
    st.markdown("**Track cancelled courses due to weather, emergencies, or other reasons**")
    
    tab1, tab2, tab3, tab4 = st.tabs(["🚫 Cancel Existing Activity", "📝 Manual Entry", "📊 View Cancellations", "📈 Analytics"])
    
    with tab1:
        st.markdown("### Cancel a Scheduled Activity")
        st.info("💡 **Smart Cancellation:** Select an existing scheduled activity. Tech time spent on setup will still count toward team hours!")
        
        # Date range to find activities
        col1, col2 = st.columns(2)
        with col1:
            search_start = st.date_input("Search From", value=datetime.now().date() - timedelta(days=7), key="cancel_search_start")
        with col2:
            search_end = st.date_input("Search To", value=datetime.now().date() + timedelta(days=30), key="cancel_search_end")
        
        # Get activities that can be cancelled
        available_activities = get_active_activities_for_cancellation(search_start, search_end)
        
        if not available_activities.empty:
            st.markdown(f"**Found {len(available_activities)} scheduled activities**")
            
            # Display activities in nice format
            display_df = available_activities[['id', 'date', 'course', 'activity_type', 'hours', 'students_trained', 'personnel', 'room_number', 'time_start']].copy()
            display_df.columns = ['ID', 'Date', 'Course', 'Activity Type', 'Hours', 'Students', 'Personnel', 'Room', 'Start Time']
            
            st.dataframe(display_df, use_container_width=True, height=300)
            
            st.markdown("---")
            st.markdown("### Cancel Selected Activity")
            
            col1, col2 = st.columns(2)
            
            with col1:
                activity_to_cancel = st.number_input("Activity ID to Cancel", min_value=1, value=int(display_df['ID'].iloc[0]) if len(display_df) > 0 else 1, step=1, key="activity_to_cancel")
                
                cancel_reason = st.selectbox(
                    "Cancellation Reason",
                    ["Weather (Snow/Ice)", "Weather (Other)", "Instructor Unavailable", "Room Unavailable", 
                     "Equipment Failure", "Low Enrollment", "Emergency", "Facility Issue", "Student Request", "Other"],
                    key="smart_cancel_reason"
                )
                
                tech_time = st.number_input("Tech Time Already Spent (hours)", min_value=0.0, max_value=24.0, value=0.0, step=0.5, key="tech_time", 
                                           help="Tech time spent on setup will still count toward team hours even though course was cancelled")
            
            with col2:
                cancel_rescheduled = st.checkbox("Will be rescheduled?", key="smart_cancel_rescheduled")
                
                if cancel_rescheduled:
                    reschedule_date = st.date_input("Reschedule Date", key="smart_reschedule_date")
                else:
                    reschedule_date = None
                
                personnel = get_personnel()
                personnel_names = sorted(personnel['name'].tolist()) if not personnel.empty else []
                cancel_logged_by = st.selectbox("Logged By", personnel_names, key="smart_cancel_logged_by") if personnel_names else None
            
            cancel_notes = st.text_area(
                "Cancellation Notes",
                placeholder="E.g., VCU closed due to winter storm. Room was already set up with SimMan 3G. Students notified via email.",
                key="smart_cancel_notes",
                height=100
            )
            
            if st.button("🚫 CANCEL THIS ACTIVITY", type="primary", use_container_width=True):
                if cancel_reason and cancel_logged_by:
                    success = cancel_existing_activity(
                        activity_id=activity_to_cancel,
                        reason=cancel_reason,
                        notes=cancel_notes,
                        tech_time_spent=tech_time,
                        rescheduled=1 if cancel_rescheduled else 0,
                        reschedule_date=reschedule_date if cancel_rescheduled else None,
                        created_by=cancel_logged_by
                    )
                    
                    if success:
                        st.success(f"✅ Activity #{activity_to_cancel} has been CANCELLED!")
                        if tech_time > 0:
                            st.info(f"💡 Tech time spent ({tech_time}h) is recorded and counts toward team hours")
                        st.warning("⚠️ This activity will NO LONGER count in Quick Stats or Dashboard metrics")
                        st.rerun()
                    else:
                        st.error(f"❌ Activity ID {activity_to_cancel} not found or already cancelled!")
                else:
                    st.error("Please select cancellation reason and who is logging this")
        
        else:
            st.info("No scheduled activities found in this date range. Try expanding the search dates or use Manual Entry tab.")
    
    with tab2:
        st.markdown("### Manual Cancellation Entry")
        st.info("💡 Use this for activities not yet entered in the system")
        
        col1, col2 = st.columns(2)
        
        with col1:
            cancel_date = st.date_input("Cancellation Date", value=datetime.now().date(), key="cancel_date")
            
            courses = get_courses()
            course_names = [""] + sorted(courses['name'].tolist()) if not courses.empty else [""]
            cancel_course = st.selectbox("Course (Optional)", course_names, key="cancel_course")
            
            cancel_time = st.time_input("Scheduled Time", key="cancel_time")
            cancel_duration = st.number_input("Scheduled Duration (hours)", min_value=0.0, max_value=24.0, value=2.0, step=0.5, key="cancel_duration")
        
        with col2:
            cancel_reason = st.selectbox(
                "Reason for Cancellation",
                ["Weather (Snow/Ice)", "Weather (Other)", "Instructor Unavailable", "Room Unavailable", 
                 "Equipment Failure", "Low Enrollment", "Emergency", "Facility Issue", "Student Request", "Other"],
                key="cancel_reason"
            )
            
            cancel_students = st.number_input("Impacted Students", min_value=0, value=0, step=1, key="cancel_students")
            
            manual_tech_time = st.number_input("Tech Time Spent (hours)", min_value=0.0, max_value=24.0, value=0.0, step=0.5, key="manual_tech_time")
            
            cancel_rescheduled = st.checkbox("Has been rescheduled?", key="cancel_rescheduled")
            
            if cancel_rescheduled:
                reschedule_date = st.date_input("Reschedule Date", key="reschedule_date")
            else:
                reschedule_date = None
        
        cancel_notes = st.text_area(
            "Notes / Additional Details",
            placeholder="E.g., VCU closed due to winter storm, all operations cancelled",
            key="cancel_notes"
        )
        
        personnel = get_personnel()
        personnel_names = sorted(personnel['name'].tolist()) if not personnel.empty else []
        cancel_created_by = st.selectbox("Logged By", personnel_names, key="cancel_created_by") if personnel_names else None
        
        if st.button("💾 Log Manual Cancellation", type="primary"):
            if cancel_reason and cancel_created_by:
                add_cancellation(
                    date=cancel_date,
                    course=cancel_course if cancel_course else None,
                    scheduled_time=cancel_time.strftime("%H:%M"),
                    scheduled_duration=cancel_duration,
                    reason=cancel_reason,
                    notes=cancel_notes,
                    impacted_students=cancel_students,
                    rescheduled=1 if cancel_rescheduled else 0,
                    reschedule_date=reschedule_date if cancel_rescheduled else None,
                    created_by=cancel_created_by,
                    tech_time_spent=manual_tech_time
                )
                st.success(f"✅ Manual cancellation logged for {cancel_date}")
                st.rerun()
            else:
                st.error("Please select a reason for cancellation and who is logging this")
    
    with tab3:
        st.markdown("### Cancellation History")
        
        col1, col2 = st.columns(2)
        with col1:
            view_start = st.date_input("From Date", value=datetime.now().date() - timedelta(days=90), key="cancel_view_start")
        with col2:
            view_end = st.date_input("To Date", value=datetime.now().date(), key="cancel_view_end")
        
        cancellations = get_cancellations(view_start, view_end)
        
        if not cancellations.empty:
            # Summary metrics
            col1, col2, col3, col4, col5 = st.columns(5)
            
            with col1:
                st.metric("Total Cancellations", len(cancellations))
            with col2:
                st.metric("Students Impacted", int(cancellations['impacted_students'].sum()))
            with col3:
                rescheduled_count = cancellations['rescheduled'].sum()
                st.metric("Rescheduled", int(rescheduled_count))
            with col4:
                total_hours_lost = cancellations['scheduled_duration'].sum()
                st.metric("Hours Lost", f"{total_hours_lost:.1f}")
            with col5:
                total_tech_time = cancellations['tech_time_spent'].sum() if 'tech_time_spent' in cancellations.columns else 0
                st.metric("Tech Hours", f"{total_tech_time:.1f}")
            
            st.markdown("---")
            
            # Filter by reason
            reason_filter = st.selectbox(
                "Filter by Reason",
                ["All"] + cancellations['reason'].unique().tolist(),
                key="cancel_reason_filter"
            )
            
            if reason_filter != "All":
                filtered_cancellations = cancellations[cancellations['reason'] == reason_filter]
            else:
                filtered_cancellations = cancellations
            
            # Display table with tech time
            st.markdown(f"**Showing {len(filtered_cancellations)} cancellations**")
            
            columns_to_show = ['id', 'date', 'course', 'scheduled_time', 'scheduled_duration', 
                              'reason', 'impacted_students', 'rescheduled', 'tech_time_spent', 'created_by']
            
            # Only show columns that exist
            columns_to_show = [col for col in columns_to_show if col in filtered_cancellations.columns]
            
            display_cancellations = filtered_cancellations[columns_to_show].copy()
            
            if 'rescheduled' in display_cancellations.columns:
                display_cancellations['rescheduled'] = display_cancellations['rescheduled'].apply(lambda x: '✅' if x == 1 else '❌')
            
            st.dataframe(display_cancellations, use_container_width=True, height=400)
            
            # Edit/Delete section
            st.markdown("---")
            st.markdown("### Edit or Delete Cancellation")
            
            cancel_id_to_edit = st.number_input("Cancellation ID to Edit", min_value=1, value=1, step=1, key="cancel_id_edit")
            
            col1, col2 = st.columns(2)
            
            with col1:
                if st.button("✏️ Mark as Rescheduled", key="mark_rescheduled"):
                    update_cancellation(cancel_id_to_edit, rescheduled=1)
                    st.success("✅ Marked as rescheduled")
                    st.rerun()
            
            with col2:
                if st.button("🗑️ Delete Cancellation", key="delete_cancel"):
                    delete_cancellation(cancel_id_to_edit)
                    st.success("✅ Cancellation deleted")
                    st.rerun()
            
            # CSV Export
            st.markdown("---")
            csv = filtered_cancellations.to_csv(index=False)
            st.download_button(
                label="📥 Download Cancellations CSV",
                data=csv,
                file_name=f"cancellations_{view_start}_to_{view_end}.csv",
                mime="text/csv"
            )
        
        else:
            st.info("No cancellations recorded for this period.")
    
    with tab4:
        st.markdown("### Cancellation Analytics")
        
        year_select = st.selectbox("Select Year", [2024, 2025, 2026], index=1, key="cancel_year")
        
        year_start = datetime(year_select, 1, 1).date()
        year_end = datetime(year_select, 12, 31).date()
        
        year_cancellations = get_cancellations(year_start, year_end)
        
        if not year_cancellations.empty:
            col1, col2 = st.columns(2)
            
            with col1:
                # Cancellations by reason
                reason_counts = year_cancellations['reason'].value_counts()
                
                fig = px.pie(
                    values=reason_counts.values,
                    names=reason_counts.index,
                    title=f"Cancellations by Reason ({year_select})",
                    color_discrete_sequence=['#F8B400', '#E74C3C', '#4A90E2', '#27AE60', '#9B59B6']
                )
                st.plotly_chart(fig, use_container_width=True)
            
            with col2:
                # Monthly trend
                year_cancellations['month'] = pd.to_datetime(year_cancellations['date']).dt.month
                monthly_counts = year_cancellations.groupby('month').size()
                
                fig = px.bar(
                    x=['Jan', 'Feb', 'Mar', 'Apr', 'May', 'Jun', 'Jul', 'Aug', 'Sep', 'Oct', 'Nov', 'Dec'],
                    y=[monthly_counts.get(i, 0) for i in range(1, 13)],
                    title=f"Monthly Cancellation Trend ({year_select})",
                    labels={'x': 'Month', 'y': 'Cancellations'},
                    color_discrete_sequence=['#F8B400']
                )
                st.plotly_chart(fig, use_container_width=True)
            
            # Impact summary
            st.markdown("---")
            st.markdown("### Impact Summary")
            
            col1, col2, col3 = st.columns(3)
            
            with col1:
                st.metric("Total Cancellations", len(year_cancellations))
                st.metric("Rescheduled", int(year_cancellations['rescheduled'].sum()))
            
            with col2:
                st.metric("Students Impacted", int(year_cancellations['impacted_students'].sum()))
                st.metric("Hours Lost", f"{year_cancellations['scheduled_duration'].sum():.1f}")
            
            with col3:
                reschedule_rate = (year_cancellations['rescheduled'].sum() / len(year_cancellations) * 100) if len(year_cancellations) > 0 else 0
                st.metric("Reschedule Rate", f"{reschedule_rate:.0f}%")
                
                avg_students = year_cancellations['impacted_students'].mean()
                st.metric("Avg Students/Cancellation", f"{avg_students:.1f}")
        
        else:
            st.info(f"No cancellations recorded for {year_select}.")

#############################################
# PAGE 5.6: TIME OFF TRACKER (NEW!)
#############################################
elif page == "🏖️ Time Off":
    st.title("🏖️ Team Time Off Tracker")
    st.markdown("**Monitor vacation, sick time, and PTO to ensure your team takes needed breaks**")
    
    tab1, tab2, tab3, tab4 = st.tabs(["📝 Log Time Off", "📊 View Records", "📈 Analytics", "💰 Leave Balances"])
    
    with tab1:
        st.markdown("### Log Team Member Time Off")
        
        personnel = get_personnel()
        personnel_names = sorted(personnel['name'].tolist()) if not personnel.empty else []
        
        col1, col2 = st.columns(2)
        
        with col1:
            pto_person = st.selectbox("Team Member", personnel_names, key="pto_person")
            pto_start = st.date_input("Start Date", value=datetime.now().date(), key="pto_start")
            pto_end = st.date_input("End Date", value=datetime.now().date(), key="pto_end")
        
        with col2:
            # Get leave types from database
            leave_types_df = get_leave_types(active_only=True)
            leave_type_options = leave_types_df['leave_type_name'].tolist() if not leave_types_df.empty else ["Vacation", "Sick Leave", "Personal Day"]
            
            pto_type = st.selectbox(
                "Type of Time Off",
                leave_type_options,
                key="pto_type"
            )
            
            # Show available hours if person selected
            if pto_person:
                accruals = get_leave_accruals(pto_person)
                if not accruals.empty:
                    matching = accruals[accruals['leave_type'] == pto_type]
                    if not matching.empty:
                        available = matching.iloc[0]['hours_available']
                        st.info(f"💡 Available: {available:.1f} hours")
                    else:
                        st.warning(f"⚠️ No accrual record for {pto_type}")
                else:
                    st.warning(f"⚠️ No leave balances initialized for {pto_person}")
            
            # Calculate hours automatically
            days_diff = (pto_end - pto_start).days + 1
            default_hours = days_diff * 8
            
            pto_hours = st.number_input("Hours", min_value=0.0, value=float(default_hours), step=0.5, key="pto_hours")
            
            pto_status = st.selectbox("Status", ["Approved", "Pending", "Denied"], key="pto_status")
        
        pto_notes = st.text_area("Notes (Optional)", placeholder="Additional information...", key="pto_notes")
        
        if st.button("💾 Log Time Off", type="primary"):
            if pto_person and pto_hours > 0:
                success, remaining = add_time_off(
                    personnel=pto_person,
                    start_date=pto_start,
                    end_date=pto_end,
                    time_off_type=pto_type,
                    hours=pto_hours,
                    status=pto_status,
                    notes=pto_notes
                )
                
                if success:
                    st.success(f"✅ Time off logged for {pto_person}")
                    st.info(f"💡 {pto_hours} hours deducted from {pto_type}. Remaining: {remaining:.1f} hours")
                else:
                    st.warning(f"⚠️ Time off logged BUT insufficient {pto_type} balance!")
                    st.error(f"❌ Attempted to use {pto_hours} hours but only {remaining:.1f} available")
                    st.info("💡 Contact HR to add accrual hours or use a different leave type")
                
                st.rerun()
            else:
                st.error("Please select team member and enter hours")
    
    with tab2:
        st.markdown("### Time Off Records")
        
        col1, col2, col3 = st.columns(3)
        with col1:
            view_start_pto = st.date_input("From Date", value=datetime.now().date() - timedelta(days=90), key="pto_view_start")
        with col2:
            view_end_pto = st.date_input("To Date", value=datetime.now().date(), key="pto_view_end")
        with col3:
            personnel_filter = st.selectbox("Filter by Person", ["All"] + personnel_names, key="pto_person_filter")
        
        if personnel_filter == "All":
            time_off_records = get_time_off(view_start_pto, view_end_pto)
        else:
            time_off_records = get_time_off(view_start_pto, view_end_pto, personnel_filter)
        
        if not time_off_records.empty:
            # Summary metrics
            col1, col2, col3, col4 = st.columns(4)
            
            with col1:
                st.metric("Total Entries", len(time_off_records))
            with col2:
                st.metric("Total Hours", f"{time_off_records['hours'].sum():.1f}")
            with col3:
                st.metric("Total Days", f"{(time_off_records['hours'].sum() / 8):.1f}")
            with col4:
                unique_people = time_off_records['personnel'].nunique()
                st.metric("Team Members", unique_people)
            
            st.markdown("---")
            
            # Display table
            st.dataframe(time_off_records, use_container_width=True, height=400)
            
            # Edit/Delete section
            st.markdown("---")
            st.markdown("### ✏️ Edit or Delete Record")
            
            pto_id_to_edit = st.number_input("Record ID to Edit", min_value=1, value=1, step=1, key="pto_id_edit")
            
            # Get the record details for display
            selected_record = time_off_records[time_off_records['id'] == pto_id_to_edit]
            
            if not selected_record.empty:
                current_person = selected_record.iloc[0]['personnel']
                current_leave_type = selected_record.iloc[0]['time_off_type']
                current_hours = selected_record.iloc[0]['hours']
                current_status = selected_record.iloc[0]['status']
                
                st.info(f"📋 Selected: {current_person} - {current_leave_type} - {current_hours} hours - Status: {current_status}")
            
            col1, col2, col3 = st.columns(3)
            
            with col1:
                st.markdown("**Change Leave Type**")
                leave_types_df = get_leave_types(active_only=True)
                leave_type_options = leave_types_df['leave_type_name'].tolist() if not leave_types_df.empty else []
                new_leave_type = st.selectbox("Update Leave Type To", leave_type_options, key="pto_new_leave_type")
                
                if st.button("✏️ Change Leave Type", type="primary"):
                    update_time_off(pto_id_to_edit, time_off_type=new_leave_type)
                    st.success(f"✅ Leave type changed to {new_leave_type}")
                    st.info("💡 Balances automatically adjusted!")
                    st.rerun()
            
            with col2:
                st.markdown("**Change Status**")
                new_status = st.selectbox("Update Status To", ["Approved", "Pending", "Denied"], key="pto_new_status")
                if st.button("✏️ Update Status"):
                    update_time_off(pto_id_to_edit, status=new_status)
                    st.success("✅ Status updated")
                    st.rerun()
            
            with col3:
                st.markdown("**Delete Record**")
                st.warning("⚠️ Will restore hours to balance")
                if st.button("🗑️ Delete Record", key="delete_pto"):
                    delete_time_off(pto_id_to_edit)
                    st.success("✅ Record deleted & hours restored")
                    st.rerun()
            
            # CSV Export
            st.markdown("---")
            csv = time_off_records.to_csv(index=False)
            st.download_button(
                label="📥 Download Time Off CSV",
                data=csv,
                file_name=f"time_off_{view_start_pto}_to_{view_end_pto}.csv",
                mime="text/csv"
            )
        
        else:
            st.info("No time off records for this period.")
    
    with tab3:
        st.markdown("### Time Off Analytics")
        
        year_select_pto = st.selectbox("Select Year", [2024, 2025, 2026], index=1, key="pto_year")
        
        summary = get_time_off_summary(year_select_pto)
        
        if not summary.empty:
            # Total time off by person
            st.markdown("### Time Off by Team Member")
            
            person_totals = summary.groupby('personnel')['total_hours'].sum().sort_values(ascending=False)
            
            fig = px.bar(
                x=person_totals.values,
                y=person_totals.index,
                orientation='h',
                title=f"Total Time Off Hours ({year_select_pto})",
                labels={'x': 'Hours', 'y': 'Team Member'},
                color_discrete_sequence=['#F8B400']
            )
            st.plotly_chart(fig, use_container_width=True)
            
            # Time off by type
            col1, col2 = st.columns(2)
            
            with col1:
                type_totals = summary.groupby('time_off_type')['total_hours'].sum()
                
                fig = px.pie(
                    values=type_totals.values,
                    names=type_totals.index,
                    title="Time Off by Type",
                    color_discrete_sequence=['#F8B400', '#4A90E2', '#27AE60', '#E67E22', '#9B59B6']
                )
                st.plotly_chart(fig, use_container_width=True)
            
            with col2:
                # Average time off per person
                avg_by_person = summary.groupby('personnel')['total_hours'].sum() / 8  # Convert to days
                
                st.markdown("#### Days Off per Person")
                for person, days in avg_by_person.sort_values(ascending=False).items():
                    st.metric(person, f"{days:.1f} days")
            
            # Detailed breakdown
            st.markdown("---")
            st.markdown("### Detailed Breakdown")
            st.dataframe(summary, use_container_width=True)
            
            # Recommendations
            st.markdown("---")
            st.markdown("### 💡 Recommendations")
            
            for person in personnel_names:
                person_data = summary[summary['personnel'] == person]
                if not person_data.empty:
                    total_days = person_data['total_hours'].sum() / 8
                    if total_days < 10:  # Less than 10 days
                        st.warning(f"⚠️ **{person}**: Only {total_days:.1f} days off this year. Encourage taking more vacation!")
                    elif total_days > 30:
                        st.info(f"ℹ️ **{person}**: {total_days:.1f} days off - good work-life balance!")
                else:
                    st.error(f"🚨 **{person}**: NO time off recorded! This is a burnout risk!")
        
        else:
            st.info(f"No time off records for {year_select_pto}.")
    
    with tab4:
        st.markdown("### 💰 Leave Balance Management")
        st.info("💡 Track available leave hours for each team member by leave type")
        
        personnel = get_personnel()
        personnel_names = sorted(personnel['name'].tolist()) if not personnel.empty else []
        
        if personnel_names:
            selected_person = st.selectbox("Select Team Member", personnel_names, key="leave_balance_person")
            
            # Initialize accruals if needed
            accruals = get_leave_accruals(selected_person)
            if accruals.empty:
                st.info(f"💡 No leave balances found for {selected_person}")
                st.markdown("**Click below to create leave balance records. All types will start at 0 hours.**")
                st.markdown("**You'll then manually add hours as needed using the 'Add/Adjust' section.**")
                
                if st.button("🆕 Initialize Leave Balances (Start at 0)", key="init_balances", type="primary"):
                    initialize_leave_accruals(selected_person)
                    st.success(f"✅ Leave balances initialized for {selected_person} - All types start at 0 hours")
                    st.info("💡 Now add hours using the section below!")
                    st.rerun()
            else:
                # Display current balances
                st.markdown(f"### Current Leave Balances for {selected_person}")
                
                balance_summary = get_leave_balance_summary(selected_person)
                
                if not balance_summary.empty:
                    for idx, row in balance_summary.iterrows():
                        leave_type = row['leave_type']
                        available = row['hours_available']
                        used = row['hours_used']
                        default_annual = row['default_annual_hours'] if pd.notna(row['default_annual_hours']) else 0
                        
                        col1, col2 = st.columns([3, 1])
                        
                        with col1:
                            st.markdown(f"**{leave_type}**")
                            total = available + used
                            if total > 0:
                                percentage = (available / total) * 100
                                st.progress(percentage / 100)
                                st.caption(f"{available:.1f} available / {total:.1f} total ({percentage:.0f}%)")
                            else:
                                st.progress(0)
                                st.caption(f"{available:.1f} available / 0.0 total")
                        
                        with col2:
                            st.metric("Used", f"{used:.1f}h")
                
                # Usage summary
                st.markdown("---")
                st.markdown("### Usage This Year")
                
                total_used = balance_summary['hours_used'].sum()
                total_days = total_used / 8
                
                col1, col2 = st.columns(2)
                with col1:
                    st.metric("Total Hours Used", f"{total_used:.1f}")
                with col2:
                    st.metric("Total Days Off", f"{total_days:.1f}")
                
                # Add/Adjust accrual section
                st.markdown("---")
                st.markdown("### ➕ Manage Accrual Hours")
                
                # Tabs for Add/Adjust vs Set Exact
                adjust_tab1, adjust_tab2 = st.tabs(["➕ Add/Subtract Hours", "🎯 Set Exact Balance"])
                
                with adjust_tab1:
                    st.info("💡 Add or subtract hours from current balance")
                    
                    col1, col2, col3 = st.columns(3)
                    
                    with col1:
                        leave_types_df = get_leave_types(active_only=True)
                        leave_type_options = leave_types_df['leave_type_name'].tolist() if not leave_types_df.empty else []
                        adjust_leave_type = st.selectbox("Leave Type", leave_type_options, key="adjust_leave_type")
                        
                        # Show current balance
                        current_accrual = balance_summary[balance_summary['leave_type'] == adjust_leave_type]
                        if not current_accrual.empty:
                            current_available = current_accrual.iloc[0]['hours_available']
                            st.caption(f"Current: {current_available:.1f} hrs")
                    
                    with col2:
                        adjust_hours = st.number_input("Hours to Add/Subtract", min_value=-1000.0, max_value=1000.0, value=0.0, step=8.0, key="adjust_hours", help="Positive = add, Negative = subtract")
                    
                    with col3:
                        st.markdown("<br>", unsafe_allow_html=True)  # Spacer
                        if st.button("💾 Add/Subtract", key="update_accrual", type="primary"):
                            if adjust_hours != 0:
                                add_accrual_hours(selected_person, adjust_leave_type, adjust_hours)
                                if adjust_hours > 0:
                                    st.success(f"✅ Added {adjust_hours} hours to {adjust_leave_type}")
                                else:
                                    st.success(f"✅ Removed {abs(adjust_hours)} hours from {adjust_leave_type}")
                                st.rerun()
                            else:
                                st.error("Please enter a non-zero amount")
                    
                    st.text_input("Reason for Adjustment", placeholder="E.g., Annual reset, Manual correction, Carryover", key="adjust_reason")
                
                with adjust_tab2:
                    st.info("💡 Set exact balance amount (overwrites current balance)")
                    
                    col1, col2, col3 = st.columns(3)
                    
                    with col1:
                        leave_types_df2 = get_leave_types(active_only=True)
                        leave_type_options2 = leave_types_df2['leave_type_name'].tolist() if not leave_types_df2.empty else []
                        set_leave_type = st.selectbox("Leave Type", leave_type_options2, key="set_leave_type")
                        
                        # Show current balance
                        current_accrual2 = balance_summary[balance_summary['leave_type'] == set_leave_type]
                        if not current_accrual2.empty:
                            current_available2 = current_accrual2.iloc[0]['hours_available']
                            st.caption(f"Current: {current_available2:.1f} hrs")
                    
                    with col2:
                        exact_hours = st.number_input("Set Balance To (Exact Hours)", min_value=0.0, max_value=5000.0, value=0.0, step=8.0, key="exact_hours", help="New total balance")
                    
                    with col3:
                        st.markdown("<br>", unsafe_allow_html=True)  # Spacer
                        if st.button("🎯 Set Exact Balance", key="set_exact_balance", type="primary"):
                            set_accrual_balance(selected_person, set_leave_type, exact_hours)
                            st.success(f"✅ {set_leave_type} balance set to {exact_hours} hours")
                            st.rerun()
                    
                    st.text_input("Reason for Change", placeholder="E.g., Correcting allocation, New policy, Carryover adjustment", key="set_reason")
                
                # Recent activity
                st.markdown("---")
                st.markdown("### 📅 Recent Time Off Activity")
                
                recent_pto = get_time_off(personnel=selected_person)
                if not recent_pto.empty:
                    recent_pto = recent_pto.sort_values('start_date', ascending=False).head(10)
                    display_pto = recent_pto[['start_date', 'time_off_type', 'hours', 'status', 'notes']].copy()
                    display_pto.columns = ['Date', 'Type', 'Hours', 'Status', 'Notes']
                    st.dataframe(display_pto, use_container_width=True, height=300)
                else:
                    st.info("No time off records found")
        else:
            st.info("No personnel found. Add team members in Data Entry first.")

#############################################
# PAGE 6: EQUIPMENT STATUS
#############################################
elif page == "🛠️ Equipment":
    st.title("🛠️ Equipment Status Dashboard")
    
    equipment_df = get_equipment()
    
    if not equipment_df.empty:
        status_counts = equipment_df['status'].value_counts()
        
        col1, col2, col3, col4 = st.columns(4)
        with col1:
            st.metric("✅ Operational", status_counts.get('Operational', 0))
        with col2:
            st.metric("🔧 Maintenance", status_counts.get('Maintenance', 0))
        with col3:
            st.metric("❌ Down", status_counts.get('Down', 0))
        with col4:
            total = len(equipment_df)
            operational = status_counts.get('Operational', 0)
            uptime = (operational / total * 100) if total > 0 else 0
            st.metric("📈 Uptime", f"{uptime:.1f}%")
        
        st.markdown("---")
        
        for idx, row in equipment_df.iterrows():
            status_icon = {"Operational": "✅", "Maintenance": "🔧", "Down": "❌"}
            with st.expander(f"{status_icon.get(row['status'], '⚪')} {row['name']} - {row['status']}"):
                st.write(f"**Status:** {row['status']}")
                st.write(f"**Last Maintenance:** {row['last_maintenance'] or 'N/A'}")
                st.write(f"**Notes:** {row['notes'] or 'None'}")
    else:
        st.info("No equipment in system. Add equipment in Settings.")

#############################################
# PAGE 7: INCIDENTS
#############################################
elif page == "⚠️ Incidents":
    st.title("⚠️ Incident Management")
    
    tab1, tab2 = st.tabs(["📋 View Incidents", "🚨 Log New Incident"])
    
    with tab1:
        active = get_incidents(resolved=0)
        resolved_inc = get_incidents(resolved=1)
        
        col1, col2, col3 = st.columns(3)
        with col1:
            st.metric("🔴 Active", len(active))
        with col2:
            st.metric("✅ Resolved", len(resolved_inc))
        with col3:
            total = len(active) + len(resolved_inc)
            rate = (len(resolved_inc) / total * 100) if total > 0 else 0
            st.metric("📊 Resolution Rate", f"{rate:.0f}%")
        
        st.markdown("---")
        
        if not active.empty:
            st.markdown("### 🔴 Active Incidents")
            for idx, row in active.iterrows():
                severity_icon = {'Critical': '🔴', 'High': '🟠', 'Medium': '🟡', 'Low': '🟢'}
                with st.expander(f"{severity_icon.get(row['severity'], '⚪')} {row['incident_type']} - {row['date']}"):
                    col1, col2 = st.columns(2)
                    with col1:
                        st.write(f"**Type:** {row['incident_type']}")
                        st.write(f"**Severity:** {row['severity']}")
                        st.write(f"**Equipment:** {row['equipment'] or 'N/A'}")
                        st.write(f"**Description:** {row['description']}")
                    with col2:
                        resolution = st.text_area("Resolution", key=f"res_{row['id']}")
                        if st.button("✅ Resolve Incident", key=f"r_{row['id']}", type="primary"):
                            if resolution:
                                resolve_incident(row['id'], resolution)
                                st.success("Incident resolved!")
                                st.rerun()
                            else:
                                st.error("Please enter resolution details")
        else:
            st.success("🎉 No active incidents!")
        
        if not resolved_inc.empty:
            st.markdown("---")
            st.markdown("### ✅ Recently Resolved")
            for idx, row in resolved_inc.head(5).iterrows():
                with st.expander(f"✅ {row['incident_type']} - {row['date']}"):
                    st.write(f"**Resolution:** {row['resolution']}")
    
    with tab2:
        st.markdown("### Log New Incident")
        
        equipment_list = [""] + get_equipment()['name'].tolist()
        
        col1, col2 = st.columns(2)
        with col1:
            inc_date = st.date_input("Date", value=datetime.now())
            inc_type = st.selectbox("Incident Type", ["Equipment Failure", "Software Issue", "Network Problem",
                                             "Safety Concern", "Training Disruption", "Facility Issue", "Other"])
            inc_equipment = st.selectbox("Equipment (optional)", equipment_list)
        with col2:
            inc_severity = st.selectbox("Severity", ["Critical", "High", "Medium", "Low"])
            inc_description = st.text_area("Description *")
        
        if st.button("🚨 Log Incident", type="primary", use_container_width=True):
            if inc_description:
                add_incident(inc_date, inc_type, inc_equipment, inc_severity, inc_description)
                st.error(f"⚠️ Incident logged: {inc_type}")
                st.rerun()
            else:
                st.error("Description required")

#############################################
# PAGE 8: GOALS
#############################################
elif page == "🎯 Goals":
    st.title("🎯 Goals & Performance Targets")
    
    tab1, tab2 = st.tabs(["📊 Progress", "🎯 Set New Goal"])
    
    with tab1:
        goals_df = get_goals()
        
        if not goals_df.empty:
            for idx, row in goals_df.iterrows():
                progress = (row['current_value'] / row['target_value'] * 100) if row['target_value'] > 0 else 0
                
                col1, col2 = st.columns([4, 1])
                with col1:
                    st.markdown(f"#### {row['goal_type']} ({row['period']})")
                with col2:
                    if st.button("🗑️", key=f"delg_{row['id']}"):
                        delete_goal(row['id'])
                        st.rerun()
                
                col1, col2, col3 = st.columns([3, 1, 1])
                with col1:
                    st.progress(min(progress / 100, 1.0))
                with col2:
                    st.metric("Progress", f"{progress:.0f}%")
                with col3:
                    st.metric("Target", f"{row['target_value']:.0f}")
                
                with st.expander("Update Progress"):
                    new_val = st.number_input("Current Value", min_value=0.0,
                                             value=float(row['current_value']), key=f"g_{row['id']}")
                    if st.button("💾 Update", key=f"ug_{row['id']}"):
                        update_goal_progress(row['id'], new_val)
                        st.success("Updated!")
                        st.rerun()
                
                st.markdown("---")
        else:
            st.info("No goals set yet. Use the 'Set New Goal' tab to add your first goal!")
    
    with tab2:
        st.markdown("### Set New Performance Goal")
        
        col1, col2 = st.columns(2)
        with col1:
            goal_type = st.text_input("Goal Description *", placeholder="e.g., Total Training Hours")
            target = st.number_input("Target Value *", min_value=0.0, step=1.0)
        with col2:
            period = st.selectbox("Time Period", ["Weekly", "Monthly", "Quarterly", "Yearly"])
        
        if st.button("🎯 Set Goal", type="primary", use_container_width=True):
            if goal_type and target > 0:
                add_goal(goal_type, target, period)
                st.success(f"✅ Goal set: {goal_type}")
                st.rerun()
            else:
                st.error("Goal description and target value required")

#############################################
# PAGE 9: SETTINGS WITH EDITING
#############################################
elif page == "⚙️ Settings":
    st.title("⚙️ System Settings")
    
    tab1, tab2, tab3, tab4, tab5, tab6 = st.tabs(["👥 Personnel", "📚 Courses", "🛠️ Equipment", "📋 Activity Types", "🏷️ Leave Types", "🏢 Room Numbers"])
    
    with tab1:
        st.markdown("### Personnel Management")
        
        col1, col2 = st.columns(2)
        
        with col1:
            st.markdown("#### Add New Personnel")
            new_name = st.text_input("Name *")
            new_role = st.text_input("Role")
            
            if st.button("➕ Add Personnel", type="primary"):
                if new_name:
                    if add_personnel(new_name, new_role):
                        st.success(f"✅ Added {new_name}")
                        st.rerun()
                    else:
                        st.error("Person already exists")
                else:
                    st.error("Name required")
        
        with col2:
            st.markdown("#### Quick Stats")
            personnel_df = get_personnel(active_only=False)
            active_personnel = len(personnel_df[personnel_df['active'] == 1])
            inactive_personnel = len(personnel_df[personnel_df['active'] == 0])
            st.metric("Active", active_personnel)
            st.metric("Inactive", inactive_personnel)
        
        st.markdown("---")
        st.markdown("#### Current Personnel")
        
        personnel_df = get_personnel(active_only=False)
        
        if not personnel_df.empty:
            for idx, row in personnel_df.iterrows():
                col1, col2 = st.columns([4, 1])
                with col1:
                    status = "✅" if row['active'] else "❌"
                    st.write(f"{status} **{row['name']}** - {row['role'] or 'No role assigned'}")
                with col2:
                    if row['active']:
                        if st.button("Deactivate", key=f"deact_p_{row['id']}"):
                            toggle_personnel(row['id'], 0)
                            st.rerun()
                    else:
                        if st.button("Activate", key=f"act_p_{row['id']}"):
                            toggle_personnel(row['id'], 1)
                            st.rerun()
    
    with tab2:
        st.markdown("### Course Management")
        
        # Add new course section
        st.markdown("#### Add New Course")
        with st.form("add_course_form"):
            col1, col2 = st.columns(2)
            
            with col1:
                new_course = st.text_input("Course Name *")
            with col2:
                new_description = st.text_area("Description")
            
            submitted = st.form_submit_button("➕ Add Course", type="primary")
            if submitted:
                if new_course:
                    if add_course(new_course, new_description):
                        st.success(f"✅ Added {new_course}")
                        st.rerun()
                    else:
                        st.error("Course already exists")
                else:
                    st.error("Course name required")
        
        st.markdown("---")
        
        # Display and edit existing courses
        st.markdown("#### Current Courses")
        
        courses_df = get_courses(active_only=False)
        
        if not courses_df.empty:
            for idx, row in courses_df.iterrows():
                status = "✅" if row['active'] else "❌"
                
                with st.expander(f"{status} **{row['name']}**"):
                    # Edit mode
                    if f"edit_course_{row['id']}" in st.session_state:
                        st.markdown("**✏️ Editing Course**")
                        
                        edit_name = st.text_input("Course Name", value=row['name'], key=f"cn_{row['id']}")
                        edit_desc = st.text_area("Description", value=row['description'] or "", key=f"cd_{row['id']}")
                        
                        col1, col2 = st.columns(2)
                        with col1:
                            if st.button("💾 Save Changes", key=f"csave_{row['id']}", type="primary"):
                                update_course(row['id'], edit_name, edit_desc)
                                del st.session_state[f"edit_course_{row['id']}"]
                                st.success("Course updated!")
                                st.rerun()
                        with col2:
                            if st.button("❌ Cancel", key=f"ccancel_{row['id']}"):
                                del st.session_state[f"edit_course_{row['id']}"]
                                st.rerun()
                    
                    # View mode
                    else:
                        if pd.notna(row.get('description')) and row.get('description'):
                            st.write(f"**Description:** {row['description']}")
                        
                        col1, col2, col3 = st.columns(3)
                        with col1:
                            if st.button("✏️ Edit", key=f"cedit_btn_{row['id']}"):
                                st.session_state[f"edit_course_{row['id']}"] = True
                                st.rerun()
                        with col2:
                            if row['active']:
                                if st.button("Deactivate", key=f"deact_c_{row['id']}"):
                                    toggle_course(row['id'], 0)
                                    st.rerun()
                            else:
                                if st.button("Activate", key=f"act_c_{row['id']}"):
                                    toggle_course(row['id'], 1)
                                    st.rerun()
        else:
            st.info("No courses yet. Add your first course above!")
    
    with tab3:
        st.markdown("### Equipment Management")
        
        col1, col2 = st.columns(2)
        
        with col1:
            st.markdown("#### Add New Equipment")
            new_equipment = st.text_input("Equipment Name *")
            new_status = st.selectbox("Initial Status", ["Operational", "Maintenance", "Down"])
            new_notes = st.text_area("Notes")
            
            if st.button("➕ Add Equipment", type="primary"):
                if new_equipment:
                    if add_equipment(new_equipment, new_status, new_notes):
                        st.success(f"✅ Added {new_equipment}")
                        st.rerun()
                    else:
                        st.error("Equipment already exists")
                else:
                    st.error("Name required")
        
        with col2:
            st.markdown("#### Quick Stats")
            equipment_df = get_equipment(active_only=False)
            active_equipment = len(equipment_df[equipment_df['active'] == 1])
            inactive_equipment = len(equipment_df[equipment_df['active'] == 0])
            st.metric("Active", active_equipment)
            st.metric("Inactive", inactive_equipment)
        
        st.markdown("---")
        st.markdown("#### Current Equipment")
        
        equipment_df = get_equipment(active_only=False)
        
        if not equipment_df.empty:
            for idx, row in equipment_df.iterrows():
                col1, col2 = st.columns([4, 1])
                with col1:
                    status = "✅" if row['active'] else "❌"
                    st.write(f"{status} **{row['name']}** - {row['status']}")
                with col2:
                    if row['active']:
                        if st.button("Deactivate", key=f"deact_e_{row['id']}"):
                            toggle_equipment(row['id'], 0)
                            st.rerun()
                    else:
                        if st.button("Activate", key=f"act_e_{row['id']}"):
                            toggle_equipment(row['id'], 1)
                            st.rerun()
    
    with tab4:
        st.markdown("### Activity Type Management")
        
        # Add new activity type
        st.markdown("#### Add New Activity Type")
        with st.form("add_activity_type_form"):
            col1, col2 = st.columns(2)
            with col1:
                new_activity = st.text_input("Activity Type Name *")
            with col2:
                new_activity_desc = st.text_input("Description")
            
            submitted = st.form_submit_button("➕ Add Activity Type", type="primary")
            if submitted:
                if new_activity:
                    if add_activity_type(new_activity, new_activity_desc):
                        st.success(f"✅ Added {new_activity}")
                        st.rerun()
                    else:
                        st.error("Activity type already exists")
                else:
                    st.error("Name required")
        
        st.markdown("---")
        
        # Display and edit activity types
        st.markdown("#### Current Activity Types")
        
        types_df = get_activity_types(active_only=False)
        
        if not types_df.empty:
            for idx, row in types_df.iterrows():
                status = "✅" if row['active'] else "❌"
                
                with st.expander(f"{status} **{row['name']}**"):
                    # Edit mode
                    if f"edit_type_{row['id']}" in st.session_state:
                        st.markdown("**✏️ Editing Activity Type**")
                        
                        edit_type_name = st.text_input("Name", value=row['name'], key=f"tn_{row['id']}")
                        edit_type_desc = st.text_input("Description", value=row['description'] or "", key=f"td_{row['id']}")
                        
                        col1, col2 = st.columns(2)
                        with col1:
                            if st.button("💾 Save Changes", key=f"tsave_{row['id']}", type="primary"):
                                update_activity_type(row['id'], edit_type_name, edit_type_desc)
                                del st.session_state[f"edit_type_{row['id']}"]
                                st.success("Activity type updated!")
                                st.rerun()
                        with col2:
                            if st.button("❌ Cancel", key=f"tcancel_{row['id']}"):
                                del st.session_state[f"edit_type_{row['id']}"]
                                st.rerun()
                    
                    # View mode
                    else:
                        if pd.notna(row.get('description')) and row.get('description'):
                            st.write(f"**Description:** {row['description']}")
                        
                        col1, col2, col3 = st.columns(3)
                        with col1:
                            if st.button("✏️ Edit", key=f"tedit_btn_{row['id']}"):
                                st.session_state[f"edit_type_{row['id']}"] = True
                                st.rerun()
                        with col2:
                            if row['active']:
                                if st.button("Deactivate", key=f"deact_at_{row['id']}"):
                                    toggle_activity_type(row['id'], 0)
                                    st.rerun()
                            else:
                                if st.button("Activate", key=f"act_at_{row['id']}"):
                                    toggle_activity_type(row['id'], 1)
                                    st.rerun()
        else:
            st.info("No activity types yet.")
    
    with tab5:
        st.markdown("### 🏷️ Leave Type Management")
        st.info("💡 Control what leave types are available for time off tracking and set default annual hours")
        
        # Add new leave type
        st.markdown("#### Add New Leave Type")
        with st.form("add_leave_type_form"):
            col1, col2 = st.columns(2)
            with col1:
                new_leave = st.text_input("Leave Type Name *", placeholder="E.g., Jury Duty")
            with col2:
                new_leave_hours = st.number_input("Default Annual Hours", min_value=0.0, value=40.0, step=8.0)
            
            submitted = st.form_submit_button("➕ Add Leave Type", type="primary")
            if submitted:
                if new_leave:
                    if add_leave_type(new_leave, new_leave_hours):
                        st.success(f"✅ Added {new_leave}")
                        st.rerun()
                    else:
                        st.error("Leave type already exists")
                else:
                    st.error("Name required")
        
        st.markdown("---")
        
        # Display current leave types
        st.markdown("#### Current Leave Types")
        
        leave_types = get_leave_types(active_only=False)
        
        if not leave_types.empty:
            col1, col2 = st.columns([3, 1])
            with col1:
                st.markdown("**Active Leave Types:**")
            with col2:
                st.markdown("**Annual Hours**")
            
            for idx, row in leave_types.iterrows():
                if row['active']:
                    col1, col2, col3 = st.columns([3, 1, 1])
                    with col1:
                        st.write(f"✅ **{row['leave_type_name']}**")
                    with col2:
                        st.write(f"{row['default_annual_hours']:.0f} hrs")
                    with col3:
                        if st.button("Deactivate", key=f"deact_lt_{row['id']}"):
                            delete_leave_type(row['id'])
                            st.rerun()
            
            st.markdown("---")
            st.markdown("**Inactive Leave Types:**")
            
            inactive = leave_types[leave_types['active'] == 0]
            if not inactive.empty:
                for idx, row in inactive.iterrows():
                    col1, col2 = st.columns([3, 1])
                    with col1:
                        st.write(f"❌ {row['leave_type_name']} ({row['default_annual_hours']:.0f} hrs)")
                    with col2:
                        if st.button("Reactivate", key=f"react_lt_{row['id']}"):
                            update_leave_type(row['id'], active=1)
                            st.rerun()
            else:
                st.info("No inactive leave types")
        else:
            st.info("No leave types yet.")
    
    with tab6:
        st.markdown("### 🏢 Room Number Management")
        st.info("💡 Add or deactivate simulation rooms. Deactivated rooms are hidden from dropdowns but preserved in history.")
        
        # Add new room
        st.markdown("#### Add New Room")
        with st.form("add_room_form"):
            col1, col2 = st.columns(2)
            with col1:
                new_room = st.text_input("Room Number *", placeholder="E.g., 9-220")
            with col2:
                new_room_notes = st.text_input("Notes (Optional)", placeholder="E.g., New sim lab")
            
            submitted = st.form_submit_button("➕ Add Room", type="primary")
            if submitted:
                if new_room:
                    if add_room_number(new_room, new_room_notes):
                        st.success(f"✅ Added Room {new_room}")
                        st.rerun()
                    else:
                        st.error("Room already exists")
                else:
                    st.error("Room number required")
        
        st.markdown("---")
        
        # Display rooms
        st.markdown("#### Current Rooms")
        
        rooms = get_room_numbers(active_only=False)
        
        if not rooms.empty:
            st.markdown("**Active Rooms:**")
            
            active_rooms = rooms[rooms['active'] == 1]
            if not active_rooms.empty:
                for idx, row in active_rooms.iterrows():
                    col1, col2 = st.columns([3, 1])
                    with col1:
                        notes_display = f" - {row['notes']}" if pd.notna(row['notes']) and row['notes'] else ""
                        st.write(f"✅ **{row['room_number']}**{notes_display}")
                    with col2:
                        if st.button("Deactivate", key=f"deact_rm_{row['id']}"):
                            toggle_room_active(row['id'])
                            st.rerun()
            
            st.markdown("---")
            st.markdown("**Inactive Rooms:**")
            
            inactive_rooms = rooms[rooms['active'] == 0]
            if not inactive_rooms.empty:
                for idx, row in inactive_rooms.iterrows():
                    col1, col2 = st.columns([3, 1])
                    with col1:
                        notes_display = f" - {row['notes']}" if pd.notna(row['notes']) and row['notes'] else ""
                        st.write(f"❌ {row['room_number']}{notes_display}")
                    with col2:
                        if st.button("Reactivate", key=f"react_rm_{row['id']}"):
                            toggle_room_active(row['id'])
                            st.rerun()
            else:
                st.info("No inactive rooms")
        else:
            st.info("No rooms yet.")

# Initialize database
init_db()
